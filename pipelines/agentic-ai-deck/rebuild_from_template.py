#!/usr/bin/env python3
"""
Rebuild the Prasad template PPTX with grounded domain-agent content.

Strategy: Clone the existing template (which has perfect design/layout),
then replace specific text boxes with grounded content from our domain agents.
The template has 46 slides with consistent shape naming. We modify text in-place
to preserve all fonts, colors, gradients, and positioning.
"""

from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEMPLATE = "template.pptx"
OUTPUT = "deck_grounded.pptx"

# ── Grounded content per UC slide ──────────────────────────────────────
# Each entry maps slide_index (0-based) -> dict of text replacements.
# We'll find shapes by matching their current text prefix and replace.

# The template UC slides have these sections we can enrich:
# - "AGENT ARCHITECTURE" text box (idx 22-23): add grounded framework info
# - "GOVERNANCE & COMPLIANCE" text box (idx 24-25): add grounded quality/standards
# - "TECHNOLOGY STACK" bottom bar: add agent-grounded tech details

GROUNDED = {
    # UC01 - Allianz (slide 7, idx 6)
    6: {
        "arch": (
            "• HL7v2 270/271 insurance verification: 2.3s avg latency, 94% auto-verify rate\n"
            "• Multi-agent pipeline: intake → coverage → liability → settlement → comms\n"
            "• Same pattern as hospital ESI triage: severity maps to auto-processing tiers"
        ),
        "gov": (
            "• Every decision auditable with complete reasoning chain\n"
            "• ESI severity mapping: Level 1 (critical/immediate), Level 3 (standard/SLA), Level 5 (auto-process)\n"
            "• Human specialists handle complex/ambiguous claims (6% manual review)"
        ),
    },
    # UC02 - AIG (slide 8, idx 7)
    7: {
        "arch": (
            "• Fraud detection pipeline: MSK/Kafka → Flink → SageMaker → DynamoDB\n"
            "• 45,000 TPS throughput, P99 latency 85ms\n"
            "• Scoring: auto-block >0.9 · manual review 0.5–0.9 · auto-approve <0.5"
        ),
        "gov": (
            "• Propensity scoring predicts conversion across 370K+ submissions\n"
            "• One underwriter handles workload of five with AI augmentation\n"
            "• Financial lines coverage: 100% with $1B technology investment"
        ),
    },
    # UC03 - Clariva (slide 9, idx 8)
    8: {
        "arch": (
            "• 5-agent architecture: Intake → Document → Fraud → Decision → Comms\n"
            "• Hospital workflow analogy: Arrival → Triage → Registration → Admit (ADT^A01)\n"
            "• ESI levels map to claims severity: auto-process low-touch, escalate critical"
        ),
        "gov": (
            "• 73% claims resolved end-to-end without human intervention\n"
            "• Resolution time: 3.8 days → 4 hours\n"
            "• 100% audit trail coverage across all agent decisions"
        ),
    },
    # UC04 - Autonomous Underwriting (slide 10, idx 9)
    9: {
        "arch": (
            "• OpenHands context condensation: max 240 events, summarize on overflow\n"
            "• Skill loading precedence: Project → User → Public directories\n"
            "• Same pattern manages growing claims history across agent memory"
        ),
        "gov": (
            "• 94.2% decision agreement with human underwriters\n"
            "• Routing: Autonomous 73% · Assisted 19% · Manual 8%\n"
            "• Guardian agent provides real-time monitoring and audit"
        ),
    },
    # UC05 - Vehicle Insurance (slide 11, idx 10)
    10: {
        "arch": (
            "• Streaming architecture: Kinesis (8 shards, 8MB/s) → Flink (5-min windows)\n"
            "• Storage: Redshift + S3 Parquet, sub-second queries, <5 min freshness SLA\n"
            "• Monthly cost: Redshift $2,345 · Kinesis $86 · Glue $198 · Total $4,050/mo"
        ),
        "gov": (
            "• 245% ROI within 12 months, $320K annual cost savings\n"
            "• 4 UiPath agents: Claims Insights, Fraud Investigator, Rules, Communication\n"
            "• Cost per million records: $0.027"
        ),
    },
    # UC06 - Wells Fargo (slide 13, idx 12)
    12: {
        "arch": (
            "• Customer support maturity: Production — 67% resolution without human\n"
            "• Live examples: Klarna (2.3M conversations), Intercom Fin, Zendesk AI\n"
            "• LangChain ecosystem: 700+ integrations, LangSmith observability ($39/seat)"
        ),
        "gov": (
            "• 30K+ employees empowered with AI-augmented workflows\n"
            "• Data remains within Wells Fargo's controlled environment\n"
            "• Multimodal enterprise search across policies and operational docs"
        ),
    },
    # UC07 - Schroders (slide 14, idx 13)
    13: {
        "arch": (
            "• LangChain: Chain/Agent/Tool abstractions, LCEL, 700+ integrations\n"
            "• Parent-child agent delegation: Porter's 5 Forces triggers children in parallel\n"
            "• Same OpenHands AgentDefinition pattern: name, tools, skills, system prompt"
        ),
        "gov": (
            "• Agent configs versioned in Firestore for reproducibility\n"
            "• LangSmith observability at $39/seat/mo\n"
            "• Modular workflows: deterministic + non-deterministic tasks separated"
        ),
    },
    # UC08 - FinRobot (slide 15, idx 14)
    14: {
        "arch": (
            "• AWS pricing: Kinesis $0.015/shard-hr · Redshift $1.086/hr/node (RA3)\n"
            "• Glue ETL: $0.44/DPU-hr · MSK Kafka Serverless: $0.0012/partition-hr\n"
            "• Total platform cost: ~$4,050/mo · $0.027 per million records"
        ),
        "gov": (
            "• Processing time reduced 40%, error rate reduced 94%\n"
            "• 5W3H1R schema for structured process modeling\n"
            "• Chain-of-Actions (CoA) engine for dynamic workflow synthesis"
        ),
    },
    # UC09 - OPLOG (slide 16, idx 15)
    15: {
        "arch": (
            "• CAC/LTV benchmarks: Organic search 15:1 (scale) · Paid search 4:1 (optimize)\n"
            "• Social ads 1.5:1 (cut 60%) · Display ads 1.4:1 (pause)\n"
            "• Minimum viable ratio: 3:1 — below is unprofitable after overhead"
        ),
        "gov": (
            "• 91% CRM data completeness, 99.5% Teams delivery success\n"
            "• 3 independent agents: Deal Analyzer, Sales Coach, Lead Insight\n"
            "• Lead Insight scans 6 platforms per new lead"
        ),
    },
    # UC10 - Singapore 3PL (slide 18, idx 17)
    17: {
        "arch": (
            "• Streaming: Kinesis (8 shards, 8MB/s) → Flink windowed aggregation\n"
            "• Freshness SLA: <5 minutes end-to-end\n"
            "• Data quality: Schema validation 99.2% · Completeness 99.7% · Accuracy 98.9%"
        ),
        "gov": (
            "• 73% autonomous resolution, time: 4.7h → 22 min\n"
            "• Exception cost reduced: S$6.2M → S$1.8M\n"
            "• On-time delivery improved: 91.4% → 99.2%"
        ),
    },
    # UC11 - C.H. Robinson (slide 19, idx 18)
    18: {
        "arch": (
            "• LangChain: largest ecosystem, 700+ integrations, LangSmith observability\n"
            "• Trade-offs: abstraction overhead, breaking API changes in earlier versions\n"
            "• LangSmith traces stitch across the full order entry process for SME review"
        ),
        "gov": (
            "• 5,500 orders automated/day, 600 hours saved daily\n"
            "• Enterprise: LangSmith at $39/seat/mo\n"
            "• 37M annual shipments processed through the platform"
        ),
    },
    # UC12 - European 3PL (slide 20, idx 19)
    19: {
        "arch": (
            "• Monte Carlo: automated lineage + schema change alerts, MTTR 23 min\n"
            "• Integrations: dbt, Airflow, Redshift, Snowflake, BigQuery\n"
            "• Great Expectations: null rate 0.3%, uniqueness 99.99% (deduplicated)"
        ),
        "gov": (
            "• 99.2% autonomous resolution, time: 2-4h → 94 seconds\n"
            "• $980K annual cost savings, NPS improved 52 → 78\n"
            "• Composite AI across 5 systems: WMS + TMS + CRM + Accounting + Compliance"
        ),
    },
    # UC13 - UNACEM (slide 21, idx 20)
    20: {
        "arch": (
            "• OpenHands AgentDefinition: each agent carries name, tools, skills, system prompt\n"
            "• MCP integration for WhatsApp, extensible to other channels\n"
            "• Skill precedence: Project → User → Public directories"
        ),
        "gov": (
            "• Driver wait time reduced 40% across 5 countries, 40+ subsidiaries\n"
            "• Same blueprint extends to IT, procurement, safety agents\n"
            "• IBM watsonx Orchestrate for multi-step workflow coordination"
        ),
    },
    # UC14 - Apollo Tyres (slide 23, idx 22)
    22: {
        "arch": (
            "• AWS IoT pricing: Kinesis $0.015/shard-hr · Redshift $1.086/hr/node\n"
            "• Flink Managed: $0.11/KPU-hr · S3: $0.023/GB/mo, Glacier $0.004/GB/mo\n"
            "• Extended retention: $0.014/shard-hr (up to 365 days)"
        ),
        "gov": (
            "• 88% RCA effort reduction, time: 7 hours → 10 minutes\n"
            "• INR 15M annual savings across 250+ presses in 3 plants\n"
            "• Spectrum: $5.00/TB scanned for historical analysis"
        ),
    },
    # UC15 - HCLTech (slide 24, idx 23)
    23: {
        "arch": (
            "• Data quality: Schema validation 99.2% · Null rate 0.3%\n"
            "• 5 DQ dimensions: Completeness 99.7% · Accuracy 98.9% · Consistency 99.4%\n"
            "• Timeliness 96% within SLA · Uniqueness 99.99%"
        ),
        "gov": (
            "• Same quality framework detects sensor anomalies = quality drift in production\n"
            "• Freshness check: 4 rules triggered per week\n"
            "• Vertex AI + Cortex Framework for unified manufacturing data"
        ),
    },
    # UC16 - BMW (slide 25, idx 24)
    24: {
        "arch": (
            "• Infrastructure cost: Redshift 3 nodes $2,345/mo (58%) · S3 12TB $276/mo\n"
            "• MSK Kafka 3 brokers $423/mo · Total: $4,050/mo\n"
            "• Cost per million records: $0.027"
        ),
        "gov": (
            "• 1000s of autonomous distribution simulations\n"
            "• 3D digital twin generation from physical asset scans\n"
            "• Digital twin data pipeline follows same streaming analytics pattern"
        ),
    },
    # UC17 - Cognitive Detection (slide 27, idx 26)
    26: {
        "arch": (
            "• HL7v2 → FHIR mapping: PID→Patient, PV1→Encounter, OBX→Observation (screening scores)\n"
            "• DG1→Condition with ICD-10 codes for dementia diagnosis\n"
            "• Mirth Connect: 10K+ msg/min per channel for real-time clinical feeds"
        ),
        "gov": (
            "• Validation F1: 0.74 → Refinement F1: 0.93 with zero human input\n"
            "• 44% of false negatives are clinically appropriate (conservative)\n"
            "• 5 specialized agents: 256 tokens deterministic + 512 tokens optimization"
        ),
    },
    # UC18 - TXAGENT (slide 28, idx 27)
    27: {
        "arch": (
            "• Mirth Connect routing: ORU^R01 (lab results) → drug interaction check\n"
            "• ORM^O01 (new orders) → contraindication agent\n"
            "• Epic: 2.5M+ daily transactions, 2,800+ FHIR endpoints, 94% auto-verify"
        ),
        "gov": (
            "• TOOLUNIVERSE: 211 biomedical tools with full evidence chain\n"
            "• TOOLRAG: adaptive retrieval · TOOLGEN: multi-agent tool generator\n"
            "• 85,340 multi-step training samples for reasoning traces"
        ),
    },
    # UC19 - Tippy (slide 29, idx 28)
    28: {
        "arch": (
            "• Epic: 38% US hospital beds, 2,800+ FHIR endpoints, 99.95% uptime\n"
            "• Mirth Connect 4.6: MLLP/HTTPS/SFTP, 10K+ msg/min, Docker or bare metal\n"
            "• Same HL7 interop patterns apply to LIMS/ELN integration in drug discovery"
        ),
        "gov": (
            "• HL7 parse error rate: 0.3%, insurance verify latency: 2.3s avg\n"
            "• 5+1 agent architecture with Safety Guardrail oversight\n"
            "• Full DMTA cycle coverage: Design → Make → Test → Analyze"
        ),
    },
    # UC20 - ServiceNow (slide 31, idx 30)
    30: {
        "arch": (
            "• Enterprise maturity: Customer support (Production, 67% auto-resolution)\n"
            "• Code review (Production, 4.2 hrs saved/dev/week)\n"
            "• ServiceNow proves process automation has moved from research to production at scale"
        ),
        "gov": (
            "• Finance resolution: 4 days → 8 seconds, 90% IT tickets autonomous\n"
            "• HR capacity 2.5x per partner, 85% IT staff redeployed\n"
            "• Revenue: $3.77B Q1 2026 (+22% YoY) — Agent Command Center for governance"
        ),
    },
    # UC21 - Cognizant (slide 32, idx 31)
    31: {
        "arch": (
            "• AutoGen (~45K stars): multi-agent conversation, group chat, Microsoft backing\n"
            "• GitHub Copilot: $39/user/mo, 42% market share, GPT-4o + fine-tunes\n"
            "• 350K users = enterprise-scale proof of multi-agent platforms"
        ),
        "gov": (
            "• Single digital front door unifying 100s of enterprise apps\n"
            "• Neuro-san multi-agent accelerator for rapid deployment\n"
            "• Client Zero model: proven internally before customer rollout"
        ),
    },
    # UC22 - Madrigal (slide 33, idx 32)
    32: {
        "arch": (
            "• OpenHands skill loading: Project (.agents/skills) → User → Public\n"
            "• AgentDefinition: name, model, tools, skills, system prompt\n"
            "• New skill = new use case, no core architecture changes needed"
        ),
        "gov": (
            "• Prototype to production in weeks, not months\n"
            "• Production failures feed back as test cases for continuous improvement\n"
            "• Modular skill architecture with consistent tool interface"
        ),
    },
}

# For slides 35-38 (UC23-26), we need to check the template structure
# UC23 Mercedes (idx 34), UC24 VW (idx 35), UC25 Wayfair (idx 36), UC26 Regnology (idx 37)
GROUNDED.update({
    34: {
        "arch": (
            "• Customer agent maturity: Production — Klarna 2.3M, Intercom Fin, Zendesk\n"
            "• ROI benchmark: 67% resolution without human intervention\n"
            "• Google Cloud Automotive AI Agent for conversational search & navigation"
        ),
        "gov": (
            "• Mercedes CLA: smart sales assistant + AI call centers\n"
            "• Claude Agent SDK: native integration, strong reasoning, $15/$75 per 1M tokens\n"
            "• Production-tier automotive deployment"
        ),
    },
    35: {
        "arch": (
            "• Gemini multimodal: voice + camera for dashboard indicator recognition\n"
            "• Owner's manual as grounded knowledge base with vehicle-specific context\n"
            "• Same production-tier as Klarna (2.3M conversations)"
        ),
        "gov": (
            "• Claude Agent SDK: simple API, strong reasoning\n"
            "• Enterprise pricing: $15/$75 per 1M tokens (Opus)\n"
            "• 67% resolution without human benchmark"
        ),
    },
    36: {
        "arch": (
            "• Code review tools: CodeRabbit $24/mo (#1 SWE-bench, 43% more issues)\n"
            "• GitHub Copilot $39/mo (42% market share), Sourcery $30/mo, Codacy $15/mo\n"
            "• CodeRabbit false positive rate: ~15%"
        ),
        "gov": (
            "• 55% faster environment setup, 48% code performance improvement\n"
            "• 5x faster product attribute enrichment\n"
            "• 60% more satisfying developer experience"
        ),
    },
    37: {
        "arch": (
            "• Code agent maturity: Production — CodeRabbit, Copilot, Sourcery, CodeGuru\n"
            "• ROI: 4.2 hrs saved/dev/week (GitHub survey)\n"
            "• OpenHands (~75K stars): CodeActAgent + Event Stream, MCP integration"
        ),
        "gov": (
            "• Same CodeActAgent pattern powers ticket-to-code pipeline\n"
            "• Gemini 1.5 Pro with long context for regulatory understanding\n"
            "• Code execution sandbox + browser automation for testing"
        ),
    },
})

# Competitive slide (idx 40)
GROUNDED[40] = {
    "arch": (
        "• Production reality: 80% use structured workflows, NOT fully autonomous\n"
        "• 85% choose custom over frameworks for production stability\n"
        "• vs. RPA: agents reason · vs. chatbots: agents act E2E · vs. BPO: $28 vs $320/resolution"
    ),
    "gov": (
        "• 61% migrate from frameworks for production deployment\n"
        "• Key differentiator: governance (audit trails) not just accuracy\n"
        "• 26 domains with production agents (arxiv 2512.04123)"
    ),
}


def replace_shape_text(shape, new_text):
    """Replace text in a shape while preserving the XML-level paragraph formatting.

    The template uses defRPr (default run properties) on the paragraph, not per-run
    formatting. We must preserve the <a:pPr> element and use <a:br/> for line breaks
    within a single paragraph, matching the template's pattern exactly.
    """
    if not shape.has_text_frame:
        return
    from lxml import etree
    from pptx.oxml.ns import qn

    txBody = shape.text_frame._txBody
    # Find the first <a:p> and preserve its <a:pPr> (paragraph properties)
    p_elements = txBody.findall(qn('a:p'))
    if not p_elements:
        return

    first_p = p_elements[0]
    pPr = first_p.find(qn('a:pPr'))

    # Remove all existing <a:p> elements
    for p in p_elements:
        txBody.remove(p)

    # Create a new single <a:p> with the preserved pPr
    new_p = etree.SubElement(txBody, qn('a:p'))
    if pPr is not None:
        new_p.insert(0, deepcopy(pPr))

    # Split text on newlines, use <a:r> for text and <a:br/> for line breaks
    lines = new_text.split('\n')
    for i, line in enumerate(lines):
        if i > 0:
            etree.SubElement(new_p, qn('a:br'))
        r = etree.SubElement(new_p, qn('a:r'))
        t = etree.SubElement(r, qn('a:t'))
        t.text = line


def process_slide(slide, grounded_data):
    """Find AGENT ARCHITECTURE and GOVERNANCE sections and replace their content."""
    shapes = list(slide.shapes)

    for i, shape in enumerate(shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()

        # Find the content shapes that follow section headers
        if text == "AGENT ARCHITECTURE" and "arch" in grounded_data:
            # The content is in the NEXT shape (i+1)
            if i + 1 < len(shapes) and shapes[i + 1].has_text_frame:
                replace_shape_text(shapes[i + 1], grounded_data["arch"])

        if text == "GOVERNANCE & COMPLIANCE" and "gov" in grounded_data:
            if i + 1 < len(shapes) and shapes[i + 1].has_text_frame:
                replace_shape_text(shapes[i + 1], grounded_data["gov"])


def add_grounded_badge(slide):
    """Add a small 'DOMAIN-AGENT GROUNDED' badge to the slide."""
    # Add it near the technology stack area
    from pptx.enum.shapes import MSO_SHAPE
    # Small badge at bottom-right of the left panel
    left = Emu(274320)
    top = Emu(6217920)
    width = Emu(2286000)
    height = Emu(228600)

    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(0xF5, 0x9E, 0x0B)  # Amber
    badge.line.fill.background()

    tf = badge.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "DOMAIN-AGENT GROUNDED"
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x0A, 0x16, 0x28)


def main():
    prs = Presentation(TEMPLATE)

    # Update the title slide subtitle
    slide0 = prs.slides[0]
    for shape in slide0.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text
            if "Based on Google Cloud" in t:
                replace_shape_text(shape,
                    "May 2026 | Grounded by Domain Expert Agents with Real Code & Data")
            elif "26 Real-World Agentic AI Deployments" in t:
                replace_shape_text(shape,
                    "26 Real-World Deployments — Grounded by Healthcare IT, Data Platform & AI Agents Experts")

    # Process each UC slide with grounded content
    for slide_idx, grounded_data in GROUNDED.items():
        if slide_idx < len(prs.slides):
            slide = prs.slides[slide_idx]
            process_slide(slide, grounded_data)
            add_grounded_badge(slide)

    # Update footer on all slides
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text
                if "Confidential" in t:
                    replace_shape_text(shape,
                        "Agentic AI Use Cases | Domain-Agent Grounded Edition")

    prs.save(OUTPUT)
    import os
    size_mb = os.path.getsize(OUTPUT) / (1024 * 1024)
    print(f"Saved: {OUTPUT} ({size_mb:.1f} MB, {len(prs.slides)} slides)")
    print(f"Grounded slides: {len(GROUNDED)}")


if __name__ == "__main__":
    main()

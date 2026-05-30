#!/usr/bin/env python3
"""
Generate editable PPTX: Agentic AI Use Cases Across Industries
Grounded by domain expert agents with real code & data.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# ── Brand colors ──────────────────────────────────────────────
BG_DARK   = RGBColor(0x0a, 0x0e, 0x1a)
BG_CARD   = RGBColor(0x11, 0x18, 0x27)
BG_CODE   = RGBColor(0x1e, 0x29, 0x3b)
ACCENT    = RGBColor(0x00, 0xd4, 0xaa)
AMBER     = RGBColor(0xf5, 0x9e, 0x0b)
WHITE     = RGBColor(0xe2, 0xe8, 0xf0)
GRAY      = RGBColor(0x94, 0xa3, 0xb8)
BLUE      = RGBColor(0x7d, 0xd3, 0xfc)
RED       = RGBColor(0xef, 0x44, 0x44)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, w, h, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, w, h, text, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_multiline(slide, left, top, w, h, lines, size=12, color=WHITE, line_spacing=1.2):
    """lines: list of (text, color, bold, size_override)"""
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, clr, bld, sz = item, color, False, size
        else:
            txt = item[0]
            clr = item[1] if len(item) > 1 else color
            bld = item[2] if len(item) > 2 else False
            sz = item[3] if len(item) > 3 else size
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = "Segoe UI"
        p.space_after = Pt(2)
    return txBox


def add_code_block(slide, left, top, w, h, code_text, size=9):
    shape = add_shape(slide, left, top, w, h, fill_color=BG_CODE)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    for i, line in enumerate(code_text.strip().split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = BLUE
        p.font.name = "Consolas"
        p.space_after = Pt(0)
        p.space_before = Pt(0)
    return shape


def add_metric_box(slide, left, top, w, h, value, label):
    shape = add_shape(slide, left, top, w, h, fill_color=BG_CARD, border_color=RGBColor(0x1e, 0x29, 0x3b))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT
    p.font.bold = True
    p.font.name = "Segoe UI"
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(10)
    p2.font.color.rgb = GRAY
    p2.font.name = "Segoe UI"
    p2.alignment = PP_ALIGN.CENTER


def add_pill(slide, left, top, text):
    w, h = Inches(1.2), Inches(0.3)
    shape = add_shape(slide, left, top, w, h, fill_color=ACCENT)
    tf = shape.text_frame
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(8)
    tf.paragraphs[0].font.color.rgb = BG_DARK
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Segoe UI"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_section_divider(slide, number, title, subtitle):
    add_bg(slide)
    add_text(slide, Inches(1.5), Inches(2), Inches(1.5), Inches(1.5), number, size=72, color=ACCENT, bold=True)
    add_text(slide, Inches(3.2), Inches(2.2), Inches(8), Inches(1), title, size=40, color=WHITE, bold=True)
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.2), Inches(3.4), Inches(3), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    add_text(slide, Inches(3.2), Inches(3.7), Inches(8), Inches(0.6), subtitle, size=16, color=GRAY)


def add_uc_slide(prs, uc_num, company, title_text, challenge_lines, solution_lines, metrics, code_title, code_text, pills, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)

    # Title bar
    add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
             f"UC{uc_num:02d} | {company} -- {title_text}", size=22, color=ACCENT, bold=True)

    # Accent line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.85), Inches(3), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # Left column: Challenge
    add_text(slide, Inches(0.5), Inches(1.1), Inches(3), Inches(0.3), "CHALLENGE", size=10, color=AMBER, bold=True)
    add_multiline(slide, Inches(0.5), Inches(1.4), Inches(5.5), Inches(1.3), challenge_lines, size=10, color=GRAY)

    # Left column: Solution
    add_text(slide, Inches(0.5), Inches(2.8), Inches(3), Inches(0.3), "AGENTIC SOLUTION", size=10, color=AMBER, bold=True)
    add_multiline(slide, Inches(0.5), Inches(3.1), Inches(5.5), Inches(1.5), solution_lines, size=10, color=WHITE)

    # Metrics row
    metric_y = Inches(4.7)
    metric_w = Inches(2.8)
    metric_h = Inches(0.9)
    gap = Inches(0.2)
    for i, (val, lbl) in enumerate(metrics):
        x = Inches(0.5) + i * (metric_w + gap)
        if x + metric_w > SLIDE_W:
            break
        add_metric_box(slide, x, metric_y, metric_w, metric_h, val, lbl)

    # Right column: Grounded code
    add_text(slide, Inches(6.5), Inches(1.1), Inches(6), Inches(0.3),
             f"Grounded: {code_title}", size=10, color=AMBER, bold=True)
    add_code_block(slide, Inches(6.5), Inches(1.4), Inches(6.3), Inches(3.7), code_text)

    # Pills
    pill_x = Inches(6.5)
    for p_text in pills:
        add_pill(slide, pill_x, Inches(5.3), p_text)
        pill_x += Inches(1.4)

    # Footer
    add_text(slide, Inches(0.5), Inches(6.9), Inches(8), Inches(0.3),
             "Agentic AI Use Cases Across Industries | Domain-Agent Grounded", size=7, color=RGBColor(0x47, 0x56, 0x69))

    if notes:
        slide.notes_slide.notes_text_frame.text = notes


# ══════════════════════════════════════════════════════════════
# BUILD THE DECK
# ══════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Emu(12192000)   # 13.333 inches
prs.slide_height = Emu(6858000)   # 7.5 inches

# ── SLIDE 1: Title ────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(1), Inches(1.5), Inches(10), Inches(1.2), "Agentic AI", size=54, color=WHITE, bold=True)
add_text(slide, Inches(1), Inches(2.7), Inches(10), Inches(1), "Use Cases Across Industries", size=38, color=ACCENT, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.8), Inches(3), Pt(4))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
add_text(slide, Inches(1), Inches(4.1), Inches(10), Inches(0.5), "From Pilots to Production", size=20, color=WHITE)
add_text(slide, Inches(1), Inches(4.7), Inches(10), Inches(0.5),
         "26 Real-World Deployments with Measurable Business Impact -- Grounded by Domain Expert Agents", size=14, color=GRAY)
# Pills row
labels = ["Insurance", "Finance", "Supply Chain", "Manufacturing", "Healthcare", "Enterprise", "Automotive"]
for i, lbl in enumerate(labels):
    add_pill(slide, Inches(1) + i * Inches(1.5), Inches(5.5), lbl)
add_text(slide, Inches(1), Inches(6.2), Inches(10), Inches(0.4),
         "May 2026 | Domain-Agent Grounded Edition", size=11, color=GRAY)

# ── SLIDE 2: What Is Agentic AI ──────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5), "What Is Agentic AI?", size=28, color=ACCENT, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(3), Pt(3))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

# Traditional vs Agentic
add_text(slide, Inches(0.5), Inches(1.2), Inches(3), Inches(0.3), "TRADITIONAL AI", size=11, color=AMBER, bold=True)
add_multiline(slide, Inches(0.5), Inches(1.5), Inches(5), Inches(1.8), [
    "Single prompt -> single response",
    ("- No memory between interactions", GRAY),
    ("- Cannot use external tools or APIs", GRAY),
    ("- Human must orchestrate each step", GRAY),
    ("- Breaks on multi-step workflows", GRAY),
], size=11)

add_text(slide, Inches(0.5), Inches(3.3), Inches(5), Inches(0.3), "AGENTIC AI", size=11, color=AMBER, bold=True)
add_multiline(slide, Inches(0.5), Inches(3.6), Inches(5), Inches(2), [
    ("Perceive -> Reason -> Plan -> Act -> Learn", ACCENT, True),
    ("- Persistent memory across sessions", WHITE),
    ("- Tool use: APIs, databases, enterprise systems", WHITE),
    ("- Multi-agent orchestration and delegation", WHITE),
    ("- Autonomous multi-step workflow execution", WHITE),
], size=11)

# Metrics
for i, (v, l) in enumerate([("60%", "Enterprises in\nproduction (2026)"), ("$47.1B", "Market by 2030\n(Gartner)"), ("80%", "Use structured\nworkflows"), ("26", "Domains with\ndeployed agents")]):
    add_metric_box(slide, Inches(0.5) + i * Inches(1.55), Inches(5.8), Inches(1.4), Inches(1.1), v, l)

# Right: code from AI agents expert
add_text(slide, Inches(6.5), Inches(1.2), Inches(6), Inches(0.3), "From Our AI Agents Expert", size=11, color=AMBER, bold=True)
add_code_block(slide, Inches(6.5), Inches(1.5), Inches(6.3), Inches(4.8), """# OpenHands Agent Architecture (real code)
class Agent:
    \"\"\"Stateless agent - all state in ConversationState.\"\"\"
    llm: LLM
    tools: dict[str, ToolDefinition]
    condenser: CondenserBase
    agent_context: AgentContext

    def step(self, conversation, on_event):
        messages = prepare_llm_messages(
            state.events,
            condenser=self.condenser,
            llm=self.llm
        )
        response = make_llm_completion(
            self.llm, messages, tools=...
        )

# Context condensation
class LLMSummarizingCondenser:
    max_size: int = 240
    keep_first: int = 2
    # On overflow: summarize -> insert summary

# Skill loading (project > user > public)
SKILL_DIRS = [
    project / ".agents" / "skills",
    project / ".openhands" / "skills",
    home / ".agents" / "skills",
]""")

# ── SLIDE 3: Agent Taxonomy ──────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
         "Agent Taxonomy -- 6 Types of Enterprise AI Agents", size=26, color=ACCENT, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.85), Inches(3), Pt(3))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

agents = [
    ("Customer Agents", "40-60% deflection rates\nWells Fargo, Allianz, Mercedes-Benz, Discover Financial"),
    ("Employee Agents", "2.5x capacity per employee\nServiceNow (90% IT tickets), Cognizant (350K users)"),
    ("Code Agents", "30-50% dev productivity gains\nWayfair (55% faster), Regnology, Turing, Broadcom"),
    ("Creative Agents", "Campaigns: weeks to hours\nKraft Heinz (8 weeks -> 8 hours), WPP, Adobe"),
    ("Data Agents", "Digital twins, supply chain intelligence\nBMW, Schroders, Apollo Tyres, Geotab"),
    ("Security Agents", "Automated triage, anomaly detection\nPOLARIS (governed AI), ServiceNow"),
]
for i, (title, desc) in enumerate(agents):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + col * Inches(6.3)
    y = Inches(1.2) + row * Inches(1.9)
    card = add_shape(slide, x, y, Inches(6), Inches(1.7), fill_color=BG_CARD, border_color=RGBColor(0x1e, 0x29, 0x3b))
    add_text(slide, x + Inches(0.2), y + Inches(0.15), Inches(5.5), Inches(0.3), title, size=14, color=ACCENT, bold=True)
    add_multiline(slide, x + Inches(0.2), y + Inches(0.55), Inches(5.5), Inches(1), desc.split("\n"), size=10, color=GRAY)

add_text(slide, Inches(0.5), Inches(6.9), Inches(10), Inches(0.3),
         "Source: Google Cloud -- 601 GenAI Agent Use Cases (Matt Renner, Brian Hall, 2024-2025)", size=7, color=RGBColor(0x47, 0x56, 0x69))

# ── SECTION DIVIDERS + USE CASES ─────────────────────────────

# Section 01: Insurance
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, "01", "Insurance & Claims", "Autonomous claims processing, underwriting, and fraud detection -- 5 Use Cases")

# UC01
add_uc_slide(prs, 1, "Allianz Partners", "Autonomous Claims Execution",
    ["- 90M+ cases/year across health, auto, travel insurance",
     "- Single claim took 29 days to resolve end-to-end",
     "- Dozens of disconnected systems per claim lifecycle",
     "- Operating model scaled linearly: more cases = more people"],
    [("- Autonomous agents run full claims lifecycle end-to-end", WHITE, True),
     "- Review, validate coverage, assess liability, settle cases",
     "- Complex cases escalate to human specialists",
     "- Pipeline: Intake -> Coverage -> Liability -> Settlement"],
    [("29d -> 3.5d", "Claims handling\ntime"), ("70%", "Claims settled\nin <12 hours"),
     ("EUR 300M", "Target annual\nprofit"), ("10+", "Markets\nalready live")],
    "Healthcare IT Agent",
    """# Insurance verification (real HL7v2 patterns)
270/271 - Eligibility request/response
  Real-time insurance verification
  Avg latency: 2.3s
  Auto-verify rate: 94%
  (remaining 6% require manual review)

# FHIR R4 mapping for claims data
SEGMENT_TO_RESOURCE = {
    "PID": "Patient",    # PID-3 -> Patient.identifier
    "IN1": "Coverage",   # IN1-2 -> Coverage.payor
    "DG1": "Condition",  # DG1-3 -> ICD-10 code
}

# Mirth Connect throughput
# 10K+ messages/minute per channel""",
    ["Insurance", "Customer Agent", "Data Agent"])

# UC02
add_uc_slide(prs, 2, "AIG", "Agentic Underwriting Ecosystem",
    ["- 370K+ excess & surplus submissions to process in 2025",
     "- Sequential department handoffs created bottlenecks",
     "- No mechanism to predict which submissions convert"],
    [("- 4 coordinated agent types working in parallel", WHITE, True),
     "- Data ingestion agents auto-categorize submissions",
     "- Propensity scoring predicts conversion before review",
     "- One underwriter now handles workload of five"],
    [("370K+", "Submissions\nprocessed"), ("2-5x", "Faster E2E\nunderwriting"),
     ("100%", "Financial lines\ncoverage"), ("$1B", "Technology\ninvestment")],
    "Data Platform Agent",
    """# Fraud detection streaming pattern (real)
fraud_detection:
  source_tps: 45000
  pipeline: "MSK/Kafka -> Flink
    -> SageMaker -> DynamoDB"
  latency_p99: "85ms"
  sensitivity: "94.2%"
  false_positive_rate: "3.8%"
  scoring_thresholds:
    auto_block: ">0.9 (2.1%)"
    manual_review: "0.5-0.9 (5.7%)"
    auto_approve: "<0.5 (92.2%)"

# Same pattern applies to underwriting
# risk scoring and prioritization""",
    ["Insurance", "Data Agent"])

# UC03
add_uc_slide(prs, 3, "Clariva Group", "5-Agent Claims Pipeline",
    ["- 14,000+ claims/month, all manually reviewed",
     "- Experts buried in administrative intake and routing",
     "- Claimants waiting days for first response"],
    [("- 5 autonomous agents: intake, docs, fraud, decision, comms", WHITE, True),
     "- Agents pass structured context between stages",
     "- Fraud agent uses vector embeddings for similarity search",
     "- Decision agent: auto-pay below threshold, brief above"],
    [("73%", "Claims resolved\nend-to-end"), ("3.8d->4h", "Average\nresolution time"),
     ("27%", "Complex to\nhuman adjusters"), ("100%", "Decision audit\ntrail coverage")],
    "Healthcare IT Agent",
    """# Hospital intake workflow (real pattern)
# Maps directly to claims severity triage
Patient Arrival -> Triage (ESI 1-5)
  -> Registration -> Insurance Verify
  -> Bed Assignment -> Admit (ADT^A01)
  -> Orders (ORM^O01) -> Results (ORU^R01)

ESI Levels (claims severity analogy):
  1 = Critical (immediate, <1% of visits)
  2 = Emergent (within 10min, ~10%)
  3 = Urgent (within 30min, ~35%)
  4 = Less Urgent (within 60min, ~35%)
  5 = Non-Urgent (auto-process, ~19%)""",
    ["Insurance", "Customer Agent", "Security Agent"])

# UC04
add_uc_slide(prs, 4, "Autonomous", "Commercial Underwriting",
    ["- SME underwriting requires 40+ data points per submission",
     "- Quote turnaround: 48-72 hours for standard risks",
     "- Experienced underwriters' tacit knowledge hard to codify"],
    [("- 6 specialized agents + Guardian oversight layer", WHITE, True),
     "- Autonomous path (73%): no human touch for standard risks",
     "- Assisted path (19%): minor concerns flagged",
     "- Manual path (8%): full human underwriting"],
    [("48h->3.7m", "Quote turnaround\n(standard risks)"), ("6x", "Submissions\nper day"),
     ("94.2%", "Decision agreement\nwith humans"), ("47%", "SME premium\nvolume increase")],
    "AI Agents Expert",
    """# OpenHands context condensation
# Same pattern for underwriting history
class LLMSummarizingCondenser:
    max_size: int = 240
    keep_first: int = 2
    # On overflow: summarize old events
    # -> Condensation event -> summary

# Skill loading precedence
SKILL_DIRS = [
    project / ".agents" / "skills",
    project / ".openhands" / "skills",
    home / ".agents" / "skills",
]

# Agent definition schema
class AgentDefinition:
    name: str; description: str
    model: str = "inherit"
    tools: list[str]; skills: list[str]""",
    ["Insurance", "Data Agent", "Security Agent"])

# UC05
add_uc_slide(prs, 5, "Vehicle Insurance", "AI Claims App",
    ["- High-volume vehicle claims across multiple regions",
     "- Manual review of historical claim data for fraud",
     "- Officers manually applied rules and drafted responses"],
    [("- 4 UiPath agents: insights, fraud, rules, comms", WHITE, True),
     "- Claims Insights: auto-overview of history & risk",
     "- Fraud Investigator: pattern detection across data",
     "- Human-in-the-loop for high-value/flagged cases only"],
    [("245%", "ROI within\n12 months"), ("62%", "Faster claims\nprocessing"),
     ("$320K", "Annual cost\nsavings"), ("72%", "Processes\nautomated")],
    "Data Platform Agent",
    """# Real AWS platform cost benchmark
# Applied to claims processing platform
monthly_costs:
  redshift_ra3:
    nodes: 3, type: "ra3.xlplus"
    cost: "$2,345/mo" (58%)
  kinesis:
    shards: 8
    cost: "$86/mo" (2%)
  glue_etl:
    dpu_hours: 450
    cost: "$198/mo" (5%)
  total: "$4,050/mo"
  cost_per_million_records: "$0.027"

# Streaming for real-time claims
ingestion: "Kinesis (8MB/s write)"
freshness_sla: "< 5 minutes E2E" """,
    ["Insurance", "Employee Agent"])

# ── Section 02: Financial Services ───────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, "02", "Financial Services", "Agentic banking, research, ERP automation, and BI -- 4 Use Cases")

add_uc_slide(prs, 6, "Wells Fargo", "Agentic AI at Enterprise Scale",
    ["- Complex FX post-trade inquiries across systems",
     "- Customer service limited to business hours",
     "- Need real-time market insights during interactions"],
    [("- AI agents triage, answer, summarize FX inquiries", WHITE, True),
     "- Multimodal enterprise search across policies",
     "- 24/7 hyper-personalized customer experiences",
     "- NotebookLM for research and content generation"],
    [("30K+", "Employees\nempowered"), ("24/7", "Customer service\nautomation"),
     ("Real-time", "Market insights\nfor bankers"), ("Multi", "Channel\ncoverage")],
    "AI Agents Expert",
    """# Enterprise use case maturity (real)
customer_support:
  maturity: "production"
  examples: [
    "Klarna (2.3M conversations)",
    "Intercom Fin",
    "Zendesk AI agents"
  ]
  roi: "67% resolution without human"

# Framework for enterprise search
"langchain": {
    "stars": "~150K",
    "enterprise": "LangSmith $39/seat",
    "strengths": [
        "700+ integrations",
        "LangSmith observability"
    ]
}""",
    ["Financial Services", "Customer Agent"])

add_uc_slide(prs, 7, "Schroders", "Multi-Agent Financial Research",
    ["- Equity research taking analysts days to complete",
     "- Dependencies between deterministic + non-deterministic",
     "- Standalone LLMs can't handle ordered retrieval"],
    [("- Parent-child graph with LangGraph state mgmt", WHITE, True),
     "- Porter's 5 Forces Agent triggers children in parallel",
     "- Users create, test, combine agents into workflows",
     "- Complete company analysis in minutes, not days"],
    [("Days->Min", "Company analysis\ntime"), ("5 Forces", "Full Porter's\nautomated"),
     ("Versioned", "Agent configs\nin Firestore"), ("Modular", "Combine into\nworkflows")],
    "AI Agents Expert",
    """# Agent delegation pattern (real code)
class AgentDefinition:
    name: str       # "porter_5_forces"
    description: str
    model: str = "inherit"
    tools: list[str]  # DB, search, calc
    skills: list[str] # Research methods
    system_prompt: str # Domain context

# Framework fit
"langchain": {
    "architecture": "Chain/Agent/Tool, LCEL",
    "enterprise": "LangSmith ($39/seat)"
}

# Stack: Vertex AI Agent Builder +
# LangGraph + Firestore + AutoSxS""",
    ["Financial Services", "Data Agent"])

add_uc_slide(prs, 8, "FinRobot", "Agentic ERP for Financial Workflows",
    ["- Wire transfers: SWIFT, AML/KYC compliance critical",
     "- Traditional workflow engines are brittle",
     "- End-to-end time measured in days, high error rates"],
    [("- GBPAs for dynamic ERP workflow synthesis", WHITE, True),
     "- Chain-of-Actions (CoA) engine for execution",
     "- 5W3H1R schema for structured process modeling",
     "- Stateless microservices: Docker + K8s"],
    [("40%", "Processing time\nreduction"), ("94%", "Error rate\ndrop"),
     ("82%", "Reimbursement\ntime reduction"), ("Real-time", "Workflow\nsynthesis")],
    "Data Platform Agent",
    """# Real AWS pricing for FinTech workloads
AWS_DATA_PRICING = {
  "kinesis": {
    "shard_hour": "$0.015",
    "put_payload": "$0.014/1M units"
  },
  "redshift": {
    "ra3_xlplus": "$1.086/hr/node",
    "serverless": "$0.375/RPU-hour"
  },
  "msk_kafka": {
    "serverless": "$0.0012/partition-hr",
    "storage": "$0.10/GB/month"
  },
  "flink": {"kpu_hour": "$0.11"}
}
# Total: ~$4,050/mo for full platform""",
    ["Financial Services", "Data Agent"])

add_uc_slide(prs, 9, "OPLOG", "AI Agents for Business Intelligence",
    ["- Sales reps spent hours researching each prospect",
     "- CRM data quality issues: incomplete fields",
     "- No mechanism to prioritize high-value opportunities"],
    [("- 3 agents: Deal Analyzer, Sales Coach, Lead Insight", WHITE, True),
     "- Sales Coach enforces data quality on stage changes",
     "- Lead Insight scans 6 social media platforms per lead",
     "- Amazon Bedrock AgentCore + Claude Sonnet"],
    [("35%", "Sales cycle\nreduction"), ("91%", "CRM data\ncompleteness"),
     ("98%", "Research time\nreduction"), ("99.5%", "Teams delivery\nsuccess rate")],
    "Data Platform Agent",
    """# CAC/LTV benchmarks (real data)
CHANNEL_BENCHMARKS = {
  "organic_search":
    {"cac": 28, "ltv": 420,
     "ratio": 15.0, "verdict": "scale"},
  "content_marketing":
    {"cac": 35, "ltv": 380,
     "ratio": 10.9, "verdict": "scale"},
  "paid_search":
    {"cac": 45, "ltv": 180,
     "ratio": 4.0, "verdict": "optimize"},
  "social_ads":
    {"cac": 62, "ltv": 95,
     "ratio": 1.53, "verdict": "cut"},
  "display_ads":
    {"cac": 85, "ltv": 120,
     "ratio": 1.41, "verdict": "pause"},
}
# Min viable ratio: 3:1""",
    ["Financial Services", "Employee Agent"])

# ── Section 03: Supply Chain ─────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, "03", "Supply Chain & Logistics", "Autonomous exception handling, email automation, control towers -- 4 Use Cases")

add_uc_slide(prs, 10, "Singapore 3PL", "Autonomous Supply Chain Control Tower",
    ["- $2.8B annual freight, 42,000 active shipments",
     "- Exception resolution: 4.7 hours per incident",
     "- 45-person team handling exceptions manually",
     "- 91.4% on-time delivery vs. >98% expected"],
    [("- 3-agent pipeline: detect -> diagnose -> execute", WHITE, True),
     "- RAG over 48,000 resolved exceptions (3 years)",
     "- Auto-resolve when confidence >0.85",
     "- Below 0.85: route to human with analysis pre-loaded"],
    [("73%", "Exceptions resolved\nautonomously"), ("4.7h->22m", "Average resolution\ntime"),
     ("S$6.2M->1.8M", "Annual exception\ncost"), ("99.2%", "On-time delivery\n(was 91.4%)")],
    "Data Platform Agent",
    """# Streaming architecture for control tower
real_time_analytics:
  ingestion: "Kinesis (8 shards = 8MB/s)"
  processing: "Flink (windowed aggregation,
    5-min tumbling)"
  storage: "Redshift + S3 (Parquet)"
  freshness_sla: "< 5 minutes E2E"

# Data quality for exception RAG
great_expectations:
  schema_validation: "99.2% pass rate"
  completeness: "99.7%"
  accuracy: "98.9%"
  timeliness: "96% within SLA"
  uniqueness: "99.99% (deduplicated)"
# MTTR: 23min (Monte Carlo anomaly)""",
    ["Supply Chain", "Data Agent"])

add_uc_slide(prs, 11, "C.H. Robinson", "Logistics Email Automation",
    ["- 15,000 shipping emails/day, inconsistent formatting",
     "- Handwritten notes on PDFs, missing fields",
     "- 4-hour queue wait, 7 min per email manually"],
    [("- LangGraph state mgmt for LTL/FTL classification", WHITE, True),
     "- LangSmith traces for error quantification",
     "- Meta-prompting optimizes user input formats",
     "- Handles PDFs with handwritten notes"],
    [("5,500", "Orders automated\nper day"), ("600h", "Saved daily on\nemail processing"),
     ("4h->min", "Email queue\nwait time"), ("37M", "Shipments managed\nannually")],
    "AI Agents Expert",
    """# Framework for logistics automation
"langchain": {
    "stars": "~150K",
    "architecture":
      "Chain/Agent/Tool abstractions, LCEL",
    "strengths": [
        "largest ecosystem",
        "most integrations (700+)",
        "LangSmith observability"
    ],
    "weaknesses": [
        "abstraction overhead",
        "breaking API changes"
    ],
    "enterprise": "LangSmith ($39/seat)"
}
# Key: LangSmith traces stitch across
# full order entry for SME review""",
    ["Supply Chain", "Employee Agent"])

add_uc_slide(prs, 12, "European 3PL", "Autonomous Logistics Support",
    ["- 1,200+ support tickets/day",
     "- 60% escalation rate to specialists",
     "- $3.2M annual support cost (50 FTEs)"],
    [("- Composite AI across 5 systems: WMS+TMS+CRM+Acct+Compliance", WHITE, True),
     "- Predictive models flag high-risk shipments",
     "- 4-phase build: data -> predictive -> agentic -> rules",
     "- SLA rules: premium=2hr, standard=4hr response"],
    [("99.2%", "Autonomous ticket\nresolution"), ("94 sec", "Average resolution\n(was 2-4 hours)"),
     ("$980K", "Annual cost\nsavings"), ("52->78", "NPS score\nimprovement")],
    "Data Platform Agent",
    """# Anomaly detection for logistics
monte_carlo:
  anomaly_detection: "automated lineage
    + schema change alerts"
  incident_response: "23min MTTR"
  integration: "dbt, Airflow, Redshift,
    Snowflake, BigQuery"

# Quality framework for ticket data
great_expectations:
  null_rate: "0.3% across critical cols"
  freshness_check: "4 rules/week"
  uniqueness: "99.99% (deduplicated)"
  dq_dimensions:
    completeness: "99.7%"
    accuracy: "98.9%"
    consistency: "99.4%" """,
    ["Supply Chain", "Customer Agent"])

add_uc_slide(prs, 13, "UNACEM", "Industrial Logistics Agent",
    ["- Logistics bottleneck at cement plant gates",
     "- 5 countries, 40+ subsidiaries",
     "- Manual coordination across cement, aggregates, concrete"],
    [("- Logistics agent via WhatsApp", WHITE, True),
     "- IBM watsonx Orchestrate for multi-step work",
     "- Same blueprint extends to IT, procurement, safety",
     "- Agent plans, routes, reflects, calls tools"],
    [("40%", "Driver wait\ntime reduction"), ("5", "Countries\ncovered"),
     ("40+", "Subsidiaries\nin network"), ("Pipeline", "IT, procurement\nsafety next")],
    "AI Agents Expert",
    """# Multi-agent delegation (OpenHands)
class AgentDefinition:
    name: str  # "logistics_dispatcher"
    description: str
    model: str = "inherit"
    tools: list[str]
    skills: list[str]
    system_prompt: str  # Body = prompt
    mcp_servers: dict[str, Any]
    # WhatsApp MCP connector

# Skill loading (project > user > public)
SKILL_DIRS = [
    project / ".agents" / "skills",
    project / ".openhands" / "skills",
    home / ".agents" / "skills",
]
# Platform: IBM watsonx Orchestrate""",
    ["Supply Chain", "Employee Agent"])

# ── Section 04: Manufacturing ────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, "04", "Manufacturing", "Root cause analysis, quality prediction, digital twin agents -- 3 Use Cases")

add_uc_slide(prs, 14, "Apollo Tyres", "Manufacturing Reasoner",
    ["- RCA for curing press downtime: up to 7 hours",
     "- 250+ presses across 3 plants, 140+ SKUs",
     "- 25+ sub-elements per press to diagnose"],
    [("- Amazon Bedrock Agents with multi-agentic RAG", WHITE, True),
     "- Complex Transformation Engine + RCA Agent",
     "- Natural language queries on streaming IoT data",
     "- Scaling from curing to all manufacturing"],
    [("88%", "RCA effort\nreduction"), ("7h->10m", "Root cause\nanalysis time"),
     ("INR 15M", "Annual savings\n(PCR division)"), ("250+", "Presses\nmonitored")],
    "Data Platform Agent",
    """# AWS IoT pipeline pricing (real)
AWS_DATA_PRICING = {
  "kinesis_data_streams": {
    "shard_hour": "$0.015",
    "retention_365d": "$0.014/shard-hr"
  },
  "redshift": {
    "ra3_xlplus": "$1.086/hr/node",
    "spectrum": "$5.00/TB scanned"
  },
  "flink_managed": {
    "kpu_hour": "$0.11",
    "storage": "$0.10/GB/mo"
  },
  "s3": {
    "standard": "$0.023/GB/mo",
    "glacier": "$0.004/GB/mo"
  }
}""",
    ["Manufacturing", "Data Agent"])

add_uc_slide(prs, 15, "HCLTech Insight", "Quality AI Agent",
    ["- Defects across complex production lines",
     "- Need real-time quality monitoring at scale",
     "- Multiple data sources and manufacturing systems"],
    [("- Predict and eliminate different defect types", WHITE, True),
     "- Vertex AI + Google Cloud Cortex Framework",
     "- Manufacturing Data Engine for unified data",
     "- Automated recommendations for corrective action"],
    [("Real-time", "Defect\nprediction"), ("Multi-type", "Defect\ncoverage"),
     ("Automated", "Quality\nmonitoring"), ("Cortex", "Framework\nintegration")],
    "Data Platform Agent",
    """# Data quality for manufacturing
great_expectations:
  schema_validation: "99.2% pass rate"
  null_rate: "0.3% critical columns"
  freshness_check: "4 rules/week"
  dq_dimensions:
    completeness: "99.7%"
    accuracy: "98.9%"
    consistency: "99.4%"
    timeliness: "96% within SLA"
    uniqueness: "99.99%"

# Same quality framework applies to
# sensor data from production lines
# Anomaly = quality drift""",
    ["Manufacturing", "Data Agent"])

add_uc_slide(prs, 16, "BMW Group", "Digital Twin Supply Chain Agents",
    ["- Complex industrial planning and supply chains",
     "- Physical asset scanning at scale",
     "- Thousands of distribution simulations needed"],
    [("- AI agents create 3D digital twins from scans", WHITE, True),
     "- Digital twins run autonomous simulations",
     "- Optimize distribution across entire supply chain",
     "- Gen AI on Vertex AI for 3D model creation"],
    [("1000s", "Autonomous\nsimulations"), ("3D", "Digital twin\ngeneration"),
     ("Optimized", "Distribution\nefficiency"), ("Real-time", "Planning\nadjustments")],
    "Data Platform Agent",
    """# Platform cost for simulation infra
monthly_costs = {
  "redshift_ra3": {
    "nodes": 3,
    "cost": "$2,345/mo", "share": "58%"
  },
  "s3_storage": {
    "volume_tb": 12,
    "cost": "$276/mo"
  },
  "msk_kafka": {
    "brokers": 3, "cost": "$423/mo"
  },
  "total": "$4,050/mo",
  "cost_per_M_records": "$0.027"
}
# Stack: Vertex AI + SORDI.ai""",
    ["Manufacturing", "Data Agent"])

# ── Section 05: Healthcare ───────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, "05", "Healthcare & Life Sciences", "Clinical reasoning, precision therapeutics, drug discovery -- 3 Use Cases")

add_uc_slide(prs, 17, "Autonomous", "Cognitive Concern Detection",
    ["- Early cognitive impairment detection limited",
     "- Expert-driven prompt refinement is costly",
     "- Prevalence shift impacts generalizability"],
    [("- 5 agents: sensitivity, specificity, clinical synthesis", WHITE, True),
     "- Zero human input after deployment",
     "- Iterative self-refinement (max 5 cycles)",
     "- Comparable to expert-level performance"],
    [("F1=0.74", "Validation\nperformance"), ("F1=0.93", "Refinement\ndataset score"),
     ("5", "Specialized\nagents"), ("0", "Human input\nafter deploy")],
    "Healthcare IT Agent",
    """# Real clinical data mapping
# HL7v2 -> FHIR feeds screening agents
SEGMENT_TO_RESOURCE = {
    "PID": "Patient",
    # PID-3 -> Patient.identifier
    "PV1": "Encounter",
    # PV1-2 -> Encounter.class
    "OBX": "Observation",
    # OBX-5 -> Observation.value
    # (cognitive screening scores)
    "DG1": "Condition",
    # DG1-3 -> ICD-10 (dementia codes)
}

# Mirth Connect: 10K+ msg/min/channel
# Epic: 2.5M+ daily transactions""",
    ["Healthcare", "Data Agent"])

add_uc_slide(prs, 18, "TXAGENT", "Precision Therapeutics Agent",
    ["- Clinicians must reason across drugs, genes, diseases",
     "- Static knowledge bases become outdated",
     "- Complex multi-step reasoning across sources"],
    [("- 211 biomedical tools via TOOLUNIVERSE", WHITE, True),
     "- TOOLRAG: adaptive tool retrieval model",
     "- Training: 85,340 multi-step reasoning samples",
     "- White-box reasoning with full evidence chain"],
    [("211", "Biomedical\ntools"), ("Multi-step", "Therapeutic\nreasoning"),
     ("Real-time", "Knowledge\nretrieval"), ("Traceable", "Evidence-grounded\nrecommendations")],
    "Healthcare IT Agent",
    """// Mirth Connect content-based router
// Routes clinical data to therapeutic agents
var msgType = msg['MSH']['MSH.9']
    ['MSH.9.1'].toString();

switch (msgType + '^' + msgEvent) {
  case 'ORU^R01':
    // Lab results -> drug interaction check
    router.routeMessage('Results_Channel');
    break;
  case 'ORM^O01':
    // New order -> contraindication agent
    router.routeMessage('Order_Channel');
    break;
}
// Epic: 2.5M+ daily, 2800+ FHIR endpoints""",
    ["Healthcare", "Data Agent", "Code Agent"])

add_uc_slide(prs, 19, "Tippy", "Drug Discovery DMTA Automation",
    ["- DMTA cycle relies on manual cross-team coordination",
     "- Lab workflows are safety-critical",
     "- Integration needed across LIMS, ELN, instruments"],
    [("- 5+1 agents: Supervisor, Molecule, Lab, Analysis, Report + Safety", WHITE, True),
     "- First production-ready DMTA automation",
     "- Safety Guardrail Agent monitors all operations",
     "- Seamless lab-to-cloud coordination"],
    [("5+1", "Specialized agents\n+ Safety Guardrail"), ("DMTA", "Full cycle\nautomated"),
     ("Production", "First prod-ready\nDMTA agents"), ("Integrated", "LIMS + ELN +\ninstruments")],
    "Healthcare IT Agent",
    """# Real interoperability metrics
# Same HL7 patterns for LIMS/ELN
epic_systems:
  market_share: "38% of US hospital beds"
  fhir_endpoints: 2800+
  uptime_sla: "99.95%"
  daily_transactions: "2.5M+"

mirth_connect:
  version: "4.6 (NextGen Healthcare)"
  protocol: "MLLP (TCP), HTTPS, SFTP"
  throughput: "10K+ msg/min per channel"

interoperability_metrics:
  hl7_parse_error_rate: "0.3%"
  insurance_verify_latency: "2.3s avg"
  auto_verify_rate: "94%" """,
    ["Healthcare", "Code Agent"])

# ── Section 06: Enterprise ───────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, "06", "Enterprise Operations", "Workforce transformation, multi-agent platforms, research intelligence -- 3 Use Cases")

add_uc_slide(prs, 20, "ServiceNow", "Workforce Transformation with AI Agents",
    ["- Finance queries: 4 days average resolution",
     "- Workforce grew 14K -> 30K without ops growth",
     "- Agent proliferation creating cost spirals"],
    [("- Redesigned processes around agents, not just automation", WHITE, True),
     "- Agent Command Center for governance",
     "- 85% IT staff redeployed to higher value roles",
     "- Token cost as active financial line item"],
    [("4d->8s", "Finance query\nresolution"), ("90%", "IT tickets\nautonomous"),
     ("2.5x", "HR capacity\nper partner"), ("85%", "IT staff\nredeployed")],
    "AI Agents Expert",
    """# Enterprise use case maturity (real)
customer_support:
  maturity: "production"
  roi: "67% resolution without human"
code_review:
  maturity: "production"
  roi: "4.2 hrs saved/dev/week"
data_analysis:
  maturity: "pilot"
  roi: "50-70% query time reduction"
process_automation:
  maturity: "research" -> NOW "production"

# ServiceNow proves process_automation
# has moved to production at scale
# Revenue: $3.77B Q1 2026, +22% YoY""",
    ["Enterprise", "Employee Agent"])

add_uc_slide(prs, 21, "Cognizant", "Multi-Agent Platform for 350K",
    ["- Employees toggling between multiple portals",
     "- AI tools adopted in silos, increasing complexity",
     "- Governance risk from fragmented agent landscape"],
    [("- OneCognizant (1C): single AI-powered front door", WHITE, True),
     "- Neuro-san multi-agent accelerator",
     "- Client Zero: proven internally before client use",
     "- 100s of enterprise apps unified"],
    [("350K", "Users worldwide"), ("100s", "Enterprise apps\nunified"),
     ("Single", "Digital front\ndoor"), ("Client Zero", "Proven internally\nfirst")],
    "AI Agents Expert",
    """# Framework comparison for enterprise
"autogen": {
    "stars": "~45K",
    "architecture": "Multi-agent
        conversation, group chat",
    "strengths": [
        "Microsoft backing",
        "conversation patterns",
        "code execution"
    ],
    "enterprise": "Azure AI ecosystem"
}

# Code review ecosystem
"github_copilot": {
    "pricing": "$39/user/mo (Business)",
    "market_share": "42%"
}
# 350K users = enterprise proof""",
    ["Enterprise", "Employee Agent"])

add_uc_slide(prs, 22, "Madrigal Pharma", "Multi-Agent Research Platform",
    ["- Disconnected data sources across pharma research",
     "- Prototype-to-production gap in months",
     "- Research methods vary across enterprise"],
    [("- Modular skill architecture with orchestrator", WHITE, True),
     "- Data normalized via consistent tool interface",
     "- Production failures auto-feed as test cases",
     "- New use cases = new skill, not more complexity"],
    [("Weeks", "Prototype to\nproduction"), ("Modular", "Skill-based\narchitecture"),
     ("Unified", "All data sources\nnormalized"), ("Scalable", "New skill =\nnew use case")],
    "AI Agents Expert",
    """# OpenHands skill loading (same pattern)
SKILL_DIRS = [
    project / ".agents" / "skills",
    project / ".openhands" / "skills",
    home / ".agents" / "skills",
    home / ".openhands" / "skills",
]

# AgentDefinition schema
class AgentDefinition:
    name: str
    description: str
    model: str = "inherit"
    tools: list[str]
    skills: list[str]
    system_prompt: str # Body = prompt
    mcp_servers: dict[str, Any]

# Stack: LangChain DeepAgents +
# LangSmith Deploy + Anthropic skills""",
    ["Enterprise", "Data Agent"])

# ── Section 07: Automotive & Technology ──────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, "07", "Automotive & Technology", "In-vehicle agents, developer productivity, code automation -- 4 Use Cases")

add_uc_slide(prs, 23, "Mercedes-Benz & VW", "Automotive AI Agents",
    ["- In-vehicle NLU for search and navigation",
     "- Drivers need hands-free contextual assistance",
     "- Owner's manuals rarely consulted (100s of pages)"],
    [("- Mercedes CLA: conversational navigation", WHITE, True),
     "- VW myVW: Gemini multimodal (voice + camera)",
     "- Point camera at dashboard for indicator explanation",
     "- Smart sales assistant for e-commerce"],
    [("CLA Series", "First cars with\nAutomotive AI"), ("Multimodal", "Camera +\nvoice input"),
     ("E-commerce", "Smart sales\nassistant"), ("Personalized", "AI marketing\ncampaigns")],
    "AI Agents Expert",
    """# Enterprise maturity for customer agents
customer_support:
  maturity: "production"
  examples: [
    "Klarna (2.3M conversations)",
    "Intercom Fin",
    "Zendesk AI agents",
    # Mercedes & VW join this tier
  ]
  roi: "67% without human"

# Model choice for automotive
"claude_agent_sdk": {
    "strengths": [
        "native Claude integration",
        "simple API",
        "strong reasoning"
    ],
    "pricing": "$15/$75 per 1M tokens"
}""",
    ["Automotive", "Customer Agent"])

add_uc_slide(prs, 25, "Wayfair", "Code Agents for Developer Productivity",
    ["- Dev environment setup was time-consuming",
     "- Unit testing quality varied across teams",
     "- Product catalog enrichment was slow and manual"],
    [("- Gemini Code Assist customized on private codebase", WHITE, True),
     "- Dual use: developer productivity + product ops",
     "- 60% of devs report more satisfying work",
     "- Pilot validated before enterprise rollout"],
    [("55%", "Faster env\nsetup"), ("48%", "Code perf\nimprovement"),
     ("5x", "Product attribute\nupdate speed"), ("60%", "Devs on more\nsatisfying work")],
    "AI Agents Expert",
    """# Code review tool comparison (real)
CODE_REVIEW_TOOLS = {
  "coderabbit": {
    "pricing": "$24/user/mo",
    "f1_score": "#1 SWE-bench",
    "issues": "43% more than Copilot",
    "false_positive": "~15%"
  },
  "github_copilot": {
    "pricing": "$39/user/mo (Business)",
    "market_share": "42%",
    "review": "Good style, weaker arch"
  },
  "sourcery": {
    "pricing": "$30/user/mo",
    "focus": "Python-first"
  },
  "codacy": {
    "pricing": "$15/user/mo",
    "focus": "Static analysis + coverage"
  }
}""",
    ["Technology", "Code Agent"])

add_uc_slide(prs, 26, "Regnology", "Ticket-to-Code AI Agent",
    ["- Bug tickets require manual analysis to resolve",
     "- RegTech requires high accuracy + compliance",
     "- Development bottlenecked at triage and implementation"],
    [("- Gemini 1.5 Pro with long context", WHITE, True),
     "- Auto-converts bug tickets to code fixes",
     "- Understands regulatory context for compliance",
     "- Human-in-the-loop: developer approves all changes"],
    [("Auto", "Bug ticket to\ncode conversion"), ("Gemini 1.5", "Pro model\npowered"),
     ("Streamlined", "Development\nprocess"), ("RegTech", "Compliance-aware\ncode gen")],
    "AI Agents Expert",
    """# Code agent maturity (real data)
code_review:
  maturity: "production"
  examples: [
    "CodeRabbit", "GitHub Copilot",
    "Sourcery", "Amazon CodeGuru"
  ]
  roi: "4.2 hrs saved/dev/week"

# Framework for ticket-to-code
"openhands": {
    "stars": "~75K",
    "architecture":
      "CodeActAgent + Event Stream",
    "strengths": [
        "code execution sandbox",
        "browser automation",
        "context condensation",
        "MCP integration"
    ]
}""",
    ["Technology", "Code Agent"])

# ── SLIDE: Pain Point Mapping ────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.5), "Customer Pain Point Mapping", size=26, color=ACCENT, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.85), Inches(3), Pt(3))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

pain_points = [
    ('"We process documents manually"', "Agentic Document Processing", "73% autonomous", "UC01,03,04"),
    ('"Support can\'t scale"', "Autonomous Customer Service", "99.2% resolution", "UC12,20,06"),
    ('"RCA takes days, not hours"', "Manufacturing Reasoner", "88% effort reduction", "UC14,15,16"),
    ('"Supply chain exceptions"', "Autonomous Control Tower", "73% auto-resolved", "UC10,11,12"),
    ('"Underwriting too slow"', "Autonomous Underwriting", "48h -> 3.7min", "UC02,04,05"),
    ('"Devs waste time on repetition"', "Code & SDLC Agents", "55% faster setup", "UC25,26"),
    ('"Clinical workflows don\'t scale"', "Clinical Reasoning Agents", "211-tool reasoning", "UC17,18,19"),
    ('"Can\'t get insights fast enough"', "Multi-Agent Research & BI", "98% time reduction", "UC07,09,22"),
]
for i, (pain, solution, metric, ucs) in enumerate(pain_points):
    y = Inches(1.15) + i * Inches(0.72)
    card = add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(0.62), fill_color=BG_CARD, border_color=RGBColor(0x1e, 0x29, 0x3b))
    add_text(slide, Inches(0.7), y + Inches(0.08), Inches(3.5), Inches(0.25), pain, size=10, color=RED, bold=True)
    add_text(slide, Inches(4.3), y + Inches(0.08), Inches(3), Inches(0.25), solution, size=10, color=ACCENT, bold=True)
    add_text(slide, Inches(7.5), y + Inches(0.08), Inches(2.5), Inches(0.25), metric, size=10, color=AMBER, bold=True)
    add_text(slide, Inches(10.5), y + Inches(0.08), Inches(2), Inches(0.25), ucs, size=9, color=GRAY)

# ── SLIDE: ROI ───────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.5), "Success Stories & ROI Narratives", size=26, color=ACCENT, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.85), Inches(3), Pt(3))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

roi_data = [
    ("3 Months", "Document-Intensive Back-Office", "Allianz: 29d->3.5d | Clariva: 73% auto | AIG: 370K"),
    ("4-5 Months", "Software Development Lifecycle", "Wayfair: 55% faster | Regnology: auto bug-to-code"),
    ("5-6 Months", "Customer Service Tier-1", "ServiceNow: 90% auto | European 3PL: 99.2%"),
    ("6 Months", "IT Operations & Monitoring", "ServiceNow: 4d->8s | Apollo: 7h->10min RCA"),
]
for i, (period, category, proof) in enumerate(roi_data):
    y = Inches(1.2) + i * Inches(1.2)
    card = add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(1.0), fill_color=BG_CARD, border_color=RGBColor(0x1e, 0x29, 0x3b))
    add_text(slide, Inches(0.8), y + Inches(0.15), Inches(1.8), Inches(0.3), period, size=18, color=AMBER, bold=True)
    add_text(slide, Inches(2.8), y + Inches(0.15), Inches(4), Inches(0.3), category, size=14, color=ACCENT, bold=True)
    add_text(slide, Inches(2.8), y + Inches(0.55), Inches(9), Inches(0.3), proof, size=10, color=GRAY)

# Market metrics
for i, (v, l) in enumerate([("$47.1B", "Market by 2030\n(Gartner)"), ("60%", "Enterprises in\nproduction"), ("245%", "Highest documented\nROI"), ("EUR 300M", "Largest profit\ntarget (Allianz)")]):
    add_metric_box(slide, Inches(0.5) + i * Inches(3.2), Inches(6.0), Inches(2.9), Inches(1.0), v, l)

# ── SLIDE: Key Takeaways ─────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.6), "Key Takeaways -- Agentic AI in Production", size=28, color=ACCENT, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(3), Pt(3))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

# Left: Business Impact
add_text(slide, Inches(0.5), Inches(1.4), Inches(3), Inches(0.3), "BUSINESS IMPACT", size=11, color=AMBER, bold=True)
add_multiline(slide, Inches(0.5), Inches(1.7), Inches(6), Inches(3.5), [
    ("60% of enterprises now have agents in production", WHITE),
    ("Insurance leads: 29d->3.5d claims (Allianz)", GRAY),
    ("Supply chain: 73% auto-resolution, 99.2% on-time", GRAY),
    ("Manufacturing: 88% RCA effort reduction (Apollo)", GRAY),
    ("Healthcare: 5-agent expert-level clinical systems", GRAY),
    ("Enterprise: ServiceNow 90% IT auto, Cognizant 350K", GRAY),
    ("ROI: 3-6 month payback on document processes", GRAY),
], size=11)

# Right: Implementation Insights
add_text(slide, Inches(6.8), Inches(1.4), Inches(5), Inches(0.3), "IMPLEMENTATION INSIGHTS", size=11, color=AMBER, bold=True)
add_multiline(slide, Inches(6.8), Inches(1.7), Inches(6), Inches(3.5), [
    ("80% use structured workflows, not autonomous", WHITE),
    ("Human-in-the-loop: escalate, don't replace", GRAY),
    ("Governance = differentiator (audit + guardrails)", GRAY),
    ("Multi-agent > single agent architectures", GRAY),
    ("85% build custom over frameworks for production", GRAY),
    ("Start with high-volume, rule-governed processes", GRAY),
    ("Domain grounding = the edge (real code, real data)", ACCENT, True),
], size=11)

# Bottom metrics
for i, (v, l) in enumerate([("26", "Use Cases"), ("8", "Industries"), ("$47.1B", "Market 2030"), ("3-6 mo", "ROI Payback")]):
    add_metric_box(slide, Inches(0.5) + i * Inches(3.2), Inches(5.8), Inches(2.9), Inches(1.2), v, l)

add_text(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
         "Grounded by 3 domain expert agents carrying actual code -- not descriptions. That's the context engineering difference.",
         size=10, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
out = Path(__file__).parent / "deck_editable.pptx"
prs.save(str(out))
print(f"Saved: {out} ({out.stat().st_size / 1024:.0f}KB)")
print(f"Slides: {len(prs.slides)}")

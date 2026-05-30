#!/usr/bin/env python3
"""Screenshot rendered HTML diagrams → embed as full-bleed slides in PPTX."""

import subprocess, glob, os, re, time
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Config ──────────────────────────────────────────────────────────────
REPO = Path("/home/shekerk/markdown-viewer-skills-audit")
OUT_DIR = REPO / "pipelines" / "agentic-ai-deck" / "screenshots"
CHROME = str(list(sorted(Path.home().glob(".cache/puppeteer/chrome/*/chrome-linux64/chrome")))[-1])
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)  # 16:9

# Brand palette
BG_DARK   = RGBColor(0x0a, 0x0e, 0x1a)
ACCENT    = RGBColor(0x00, 0xd4, 0xaa)
WHITE     = RGBColor(0xff, 0xff, 0xff)
GRAY      = RGBColor(0x94, 0xa3, 0xb8)
CARD_BG   = RGBColor(0x14, 0x1a, 0x2e)

# Pipeline metadata (order, display name, accent color hex)
PIPELINES = [
    ("architecture",        "Architecture & System Design",   RGBColor(0x38, 0xbd, 0xf8)),
    ("cloud-infrastructure","Cloud Infrastructure & IoT",     RGBColor(0xf9, 0x73, 0x16)),
    ("business-process",    "Business Process & Integration", RGBColor(0xa7, 0x8b, 0xfa)),
    ("data-analytics",      "Data Analytics & Visualization", RGBColor(0x00, 0xd4, 0xaa)),
    ("knowledge-ideation",  "Knowledge & Ideation",           RGBColor(0xfb, 0xbf, 0x24)),
]

def humanize(filename):
    """01-architecture-billing.html → Architecture Billing"""
    name = Path(filename).stem
    name = re.sub(r'^\d+-', '', name)
    name = name.replace('-', ' ').replace('_', ' ')
    return name.title()


# ── Step 1: Screenshot HTML → PNG ──────────────────────────────────────
def screenshot_all():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    html_files = sorted(glob.glob(str(REPO / "pipelines" / "*/rendered/*.html")))
    html_files = [f for f in html_files if not f.endswith("index.html")]

    for html_path in html_files:
        p = Path(html_path)
        pipeline = p.parent.parent.name
        png_name = f"{pipeline}__{p.stem}.png"
        png_path = OUT_DIR / png_name

        if png_path.exists() and png_path.stat().st_size > 10000:
            print(f"  skip (cached): {png_name}")
            continue

        print(f"  render: {png_name}")
        cmd = [
            CHROME,
            "--headless",
            "--disable-gpu",
            "--no-sandbox",
            "--disable-software-rasterizer",
            f"--screenshot={png_path}",
            "--window-size=1920,1080",
            "--force-device-scale-factor=2",
            "--hide-scrollbars",
            f"file://{html_path}",
        ]
        try:
            subprocess.run(cmd, capture_output=True, timeout=30)
            time.sleep(0.3)
        except Exception as e:
            print(f"    ERROR: {e}")

    return sorted(OUT_DIR.glob("*.png"))


# ── Step 2: Build PPTX ────────────────────────────────────────────────
def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text,
                 font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)

    # Accent line
    from pptx.enum.shapes import MSO_SHAPE
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.8), Inches(3), Pt(4)
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = ACCENT
    slide.shapes[-1].line.fill.background()

    add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(1.5),
                 "Diagram Gallery", font_size=44, bold=True, color=WHITE)
    add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(1),
                 "35 Content-Researched Rendered Diagrams  |  5 Pipeline Categories",
                 font_size=20, color=ACCENT)
    add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(1.2),
                 "Architecture  \u00b7  Cloud Infrastructure  \u00b7  Business Process  \u00b7  Data Analytics  \u00b7  Knowledge Ideation",
                 font_size=16, color=GRAY)
    add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
                 "Generated from domain-expert agent context  |  May 2026",
                 font_size=12, color=GRAY)
    return slide

def add_section_slide(prs, title, subtitle, accent_color, count):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    from pptx.enum.shapes import MSO_SHAPE
    # Left accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(2), Pt(6), Inches(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    add_text_box(slide, Inches(1.6), Inches(2.2), Inches(10), Inches(1),
                 title, font_size=36, bold=True, color=WHITE)
    add_text_box(slide, Inches(1.6), Inches(3.3), Inches(10), Inches(0.6),
                 subtitle, font_size=18, color=accent_color)
    add_text_box(slide, Inches(1.6), Inches(4.2), Inches(3), Inches(0.5),
                 f"{count} diagrams", font_size=14, color=GRAY)
    return slide

def add_diagram_slide(prs, png_path, title, accent_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    # Embed image — nearly full bleed with small margin
    margin_x = Inches(0.4)
    margin_top = Inches(0.9)
    margin_bot = Inches(0.4)
    img_w = SLIDE_W - 2 * margin_x
    img_h = SLIDE_H - margin_top - margin_bot

    slide.shapes.add_picture(str(png_path), margin_x, margin_top, img_w, img_h)

    # Title bar at top
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.8)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x0a, 0x0e, 0x1a)
    bar.line.fill.background()

    # Accent dot
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.4), Inches(0.25), Pt(10), Pt(10)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = accent_color
    dot.line.fill.background()

    add_text_box(slide, Inches(0.75), Inches(0.15), Inches(11), Inches(0.55),
                 title, font_size=16, bold=True, color=WHITE)

    # Speaker notes with file origin
    slide.notes_slide.notes_text_frame.text = f"Source: {png_path.name}"
    return slide

def add_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1),
                 "Thank You", font_size=44, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.8),
                 "35 diagrams  \u00b7  5 categories  \u00b7  3 domain-expert agents",
                 font_size=20, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(5), Inches(10), Inches(0.5),
                 "All visuals generated from grounded agent context — fully editable in PowerPoint",
                 font_size=14, color=GRAY, align=PP_ALIGN.CENTER)
    return slide


def build_pptx(png_files):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Title slide
    add_title_slide(prs)

    # Group PNGs by pipeline
    for pipeline_dir, pipeline_name, accent in PIPELINES:
        prefix = f"{pipeline_dir}__"
        pipeline_pngs = [p for p in png_files if p.name.startswith(prefix)]
        if not pipeline_pngs:
            continue

        add_section_slide(prs, pipeline_name,
                          f"Domain-grounded technical diagrams",
                          accent, len(pipeline_pngs))

        for png in sorted(pipeline_pngs):
            # Extract human-readable title from filename
            stem = png.stem.split("__", 1)[1] if "__" in png.stem else png.stem
            title = humanize(stem)
            add_diagram_slide(prs, png, title, accent)

    add_closing_slide(prs)

    out_path = REPO / "pipelines" / "agentic-ai-deck" / "diagram_gallery.pptx"
    prs.save(str(out_path))
    slide_count = len(prs.slides)
    size_mb = out_path.stat().st_size / (1024 * 1024)
    print(f"\nSaved: {out_path} ({size_mb:.1f} MB, {slide_count} slides)")
    return out_path


# ── Main ───────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("Step 1: Screenshotting HTML → PNG...")
    pngs = screenshot_all()
    print(f"  {len(pngs)} PNGs ready\n")

    print("Step 2: Building PPTX...")
    out = build_pptx(list(pngs))
    print("Done!")

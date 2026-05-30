#!/usr/bin/env python3
"""Generate an editable PPTX from the architecture pipeline test outputs."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Theme Colors ──
NAVY = RGBColor(0x1E, 0x1B, 0x4B)
INDIGO = RGBColor(0x63, 0x66, 0xF1)
INDIGO_LIGHT = RGBColor(0xA5, 0xB4, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_100 = RGBColor(0xF1, 0xF5, 0xF9)
GRAY_400 = RGBColor(0x94, 0xA3, 0xB8)
GRAY_600 = RGBColor(0x47, 0x56, 0x69)
GRAY_900 = RGBColor(0x0F, 0x17, 0x2A)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)
ORANGE = RGBColor(0xF9, 0x73, 0x16)

# Layer colors
LAYER_USER = RGBColor(0x3B, 0x82, 0xF6)
LAYER_APP = RGBColor(0xD9, 0x77, 0x06)
LAYER_DATA = RGBColor(0xDB, 0x27, 0x77)
LAYER_INFRA = RGBColor(0x6B, 0x72, 0x80)
LAYER_EXT = RGBColor(0x81, 0x8C, 0xF8)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=12, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox


def add_box(slide, left, top, width, height, text, fill_color, border_color, text_color=GRAY_900, size=9, bold=True, detail=""):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    if detail:
        p2 = tf.add_paragraph()
        p2.text = detail
        p2.font.size = Pt(7)
        p2.font.color.rgb = GRAY_600
        p2.alignment = PP_ALIGN.CENTER
    return shape


def add_layer_row(slide, top, label, label_color, items, item_colors=None):
    """Add a layer label + row of component boxes."""
    add_text(slide, 0.4, top, 1.2, 0.3, label, size=8, bold=True, color=label_color)
    w = min(1.8, 7.5 / len(items))
    gap = 0.1
    start_x = 1.7
    for i, item in enumerate(items):
        name = item if isinstance(item, str) else item[0]
        detail = "" if isinstance(item, str) else item[1]
        fc = WHITE
        bc = label_color
        add_box(slide, start_x + i * (w + gap), top, w, 0.45, name, fc, bc, GRAY_900, size=8, detail=detail)


def add_metric_card(slide, left, top, title, subtitle, metrics, accent_color):
    """Add a metric card shape group."""
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(4.3), Inches(2.2))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    card.line.color.rgb = RGBColor(0x33, 0x41, 0x55)
    card.line.width = Pt(1)
    # Category
    add_text(slide, left + 0.2, top + 0.1, 3.5, 0.2, subtitle.upper(), size=7, bold=True, color=accent_color)
    # Title
    add_text(slide, left + 0.2, top + 0.35, 3.5, 0.3, title, size=16, bold=True, color=GRAY_100)
    # Metrics row
    mx = left + 0.2
    for label, value, color in metrics:
        m_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(mx), Inches(top + 0.85), Inches(1.2), Inches(0.7))
        m_bg.fill.solid()
        m_bg.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        m_bg.line.color.rgb = RGBColor(0x33, 0x41, 0x55)
        m_bg.line.width = Pt(0.5)
        add_text(slide, mx + 0.1, top + 0.9, 1.0, 0.25, value, size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, mx + 0.1, top + 1.2, 1.0, 0.2, label, size=7, color=GRAY_600, align=PP_ALIGN.CENTER)
        mx += 1.3


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ══════════════════════════════════════════════════
    # SLIDE 1: Title
    # ══════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, NAVY)
    add_text(slide, 1, 1.5, 11, 1, "Architecture Pipeline Tests", size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, 1, 2.8, 11, 0.5, "SaaS Billing Platform  |  5 Skills Tested  |  Editable Slides", size=18, color=INDIGO_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, 1, 4.0, 11, 0.4, "Skills: architecture  /  infocard  /  drawio  /  workflow-visualizer  /  slide-narrative", size=12, color=GRAY_400, align=PP_ALIGN.CENTER)
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.6), Inches(4.3), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = INDIGO
    line.line.fill.background()

    # ══════════════════════════════════════════════════
    # SLIDE 2: Architecture — Billing Platform (indigo-deep, three-column)
    # ══════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    add_text(slide, 0.4, 0.2, 12, 0.5, "Test 1: Multi-Tenant SaaS Billing Platform", size=22, bold=True, color=WHITE)
    add_text(slide, 0.4, 0.65, 6, 0.3, "Skill: architecture  |  Style: indigo-deep  |  Layout: three-column", size=10, color=GRAY_400)

    # Left sidebar
    add_text(slide, 0.3, 1.2, 1.3, 0.25, "OBSERVABILITY", size=8, bold=True, color=INDIGO_LIGHT, align=PP_ALIGN.CENTER)
    for i, item in enumerate(["Datadog APM", "Structured Logging", "Distributed Tracing", "99.99% Uptime SLA"]):
        c = INDIGO if i == 3 else GRAY_400
        add_text(slide, 0.3, 1.5 + i * 0.25, 1.3, 0.2, item, size=8, color=c, align=PP_ALIGN.CENTER)
    add_text(slide, 0.3, 2.8, 1.3, 0.25, "ALERTING", size=8, bold=True, color=INDIGO_LIGHT, align=PP_ALIGN.CENTER)
    for i, item in enumerate(["PagerDuty On-Call", "Anomaly Alerts", "Revenue Dashboards"]):
        add_text(slide, 0.3, 3.1 + i * 0.25, 1.3, 0.2, item, size=8, color=GRAY_400, align=PP_ALIGN.CENTER)

    # Main layers
    add_layer_row(slide, 1.2, "USER", LAYER_USER,
                  [("Customer Portal", "React"), ("Admin Console", "React"), ("API Clients", "REST/GraphQL")])
    add_layer_row(slide, 1.85, "APPLICATION", LAYER_APP,
                  [("API Gateway", "Kong"), ("Billing Engine", "Metering"), ("Subscription Mgr", "Plans")])
    add_layer_row(slide, 2.45, "APP (cont.)", LAYER_APP,
                  [("Invoice Generator", "PDF/Tax"), ("Stripe Webhooks", "Events"), ("Notifications", "Email/Slack")])
    add_layer_row(slide, 3.1, "DATA", LAYER_DATA,
                  [("PostgreSQL", "Subs/Invoices"), ("Redis", "Counters/Cache"), ("S3", "Invoice PDFs")])
    add_layer_row(slide, 3.7, "INFRA", LAYER_INFRA,
                  [("EKS", "Kubernetes"), ("RDS Multi-AZ", "Primary+Replica"), ("ElastiCache", "Redis Cluster"), ("CloudFront", "CDN")])
    add_layer_row(slide, 4.35, "EXTERNAL", LAYER_EXT,
                  [("Stripe", "Payments"), ("Avalara", "Tax"), ("SendGrid", "Email"), ("Auth0", "Identity")])

    # Right sidebar
    add_text(slide, 11.5, 1.2, 1.5, 0.25, "PCI COMPLIANCE", size=8, bold=True, color=INDIGO_LIGHT, align=PP_ALIGN.CENTER)
    for i, item in enumerate(["PCI DSS Level 1", "Tokenized Cards", "Encrypted at Rest", "TLS 1.3 in Transit"]):
        add_text(slide, 11.5, 1.5 + i * 0.25, 1.5, 0.2, item, size=8, color=GRAY_400, align=PP_ALIGN.CENTER)
    add_text(slide, 11.5, 2.8, 1.5, 0.25, "SOX COMPLIANCE", size=8, bold=True, color=INDIGO_LIGHT, align=PP_ALIGN.CENTER)
    for i, item in enumerate(["Revenue Recognition", "Audit Trail", "Segregation of Duties"]):
        add_text(slide, 11.5, 3.1 + i * 0.25, 1.5, 0.2, item, size=8, color=GRAY_400, align=PP_ALIGN.CENTER)
    add_text(slide, 11.5, 4.1, 1.5, 0.25, "MULTI-TENANCY", size=8, bold=True, color=INDIGO_LIGHT, align=PP_ALIGN.CENTER)
    for i, item in enumerate(["Tenant Isolation", "Row-Level Security", "<200ms P99"]):
        c = INDIGO if i == 2 else GRAY_400
        add_text(slide, 11.5, 4.4 + i * 0.25, 1.5, 0.2, item, size=8, color=c, align=PP_ALIGN.CENTER)

    # Grade box
    add_text(slide, 0.4, 5.3, 12, 0.3, "GRADE: PASS  |  Direct HTML (no fence)  |  No empty lines  |  indigo-deep style  |  three-column layout  |  All layers populated  |  Both sidebars concrete",
             size=9, color=GREEN)

    # ══════════════════════════════════════════════════
    # SLIDE 3: Architecture — URL Shortener (frost-clean, single-stack)
    # ══════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, GRAY_100)
    add_text(slide, 0.4, 0.2, 12, 0.5, "Test 2: URL Shortener Microservice", size=22, bold=True, color=GRAY_900)
    add_text(slide, 0.4, 0.65, 6, 0.3, "Skill: architecture  |  Style: frost-clean  |  Layout: single-stack", size=10, color=GRAY_600)

    layer_data = [
        ("USER", LAYER_USER, [("Web Dashboard", "Create & Manage"), ("REST API", "Programmatic"), ("Browser Redirect", "Short URL -> Target")]),
        ("APPLICATION", LAYER_APP, [("Shortener Service", "Base62 Encoding"), ("Redirect Handler", "301/302"), ("Analytics Collector", "Click Tracking")]),
        ("DATA", LAYER_DATA, [("PostgreSQL", "URL Mappings"), ("Redis", "Hot Cache"), ("ClickHouse", "Click Analytics")]),
        ("INFRA", LAYER_INFRA, [("Docker", "Runtime"), ("Nginx", "Reverse Proxy"), ("CloudFlare", "CDN/DDoS"), ("GitHub Actions", "CI/CD")]),
    ]
    for i, (label, color, items) in enumerate(layer_data):
        y = 1.2 + i * 0.75
        add_text(slide, 1.5, y, 1.5, 0.25, label, size=9, bold=True, color=color, align=PP_ALIGN.CENTER)
        w = 2.2
        gap = 0.15
        start_x = 3.0
        for j, (name, detail) in enumerate(items):
            add_box(slide, start_x + j * (w + gap), y, w, 0.5, name, WHITE, color, GRAY_900, size=9, detail=detail)

    add_text(slide, 0.4, 4.5, 12, 0.3, "GRADE: PASS  |  single-stack layout  |  frost-clean style  |  4 layers  |  No sidebars  |  Proportionate complexity",
             size=9, color=RGBColor(0x16, 0xA3, 0x4A))

    # ══════════════════════════════════════════════════
    # SLIDE 4: Info Cards — Component Metrics
    # ══════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, GRAY_900)
    add_text(slide, 0.4, 0.2, 12, 0.5, "Test 3: Component Metric Cards", size=22, bold=True, color=WHITE)
    add_text(slide, 0.4, 0.65, 6, 0.3, "Skill: infocard  |  Layout: metric-board  |  Tone: technical", size=10, color=GRAY_400)

    add_metric_card(slide, 0.4, 1.2, "Billing Engine", "Core Service",
                    [("Invoices/Day", "12K", INDIGO_LIGHT), ("P99 Latency", "45ms", GREEN), ("Uptime (90d)", "99.97%", AMBER)], INDIGO)
    add_metric_card(slide, 5.0, 1.2, "Subscription Manager", "Plan Management",
                    [("Active Subs", "8.4K", INDIGO_LIGHT), ("Renewal Rate", "94%", GREEN), ("Plan Tiers", "12", AMBER)], CYAN)
    add_metric_card(slide, 0.4, 3.7, "Stripe Webhook Handler", "Integration",
                    [("Event Types", "28", INDIGO_LIGHT), ("Avg Response", "150ms", GREEN), ("Failure Rate", "0.02%", RED)], PURPLE)
    add_metric_card(slide, 5.0, 3.7, "Invoice Generator", "Output Service",
                    [("PDFs/Month", "12K", INDIGO_LIGHT), ("Render Time", "1.8s", GREEN), ("Currencies", "6", AMBER)], ORANGE)

    add_text(slide, 0.4, 6.2, 12, 0.3, "GRADE: PASS  |  4 cards  |  Direct HTML  |  tech-blueprint tone auto-sensed  |  Metrics prominent  |  No code fence",
             size=9, color=GREEN)

    # ══════════════════════════════════════════════════
    # SLIDE 5: Draw.io — Billing Request Flow
    # ══════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text(slide, 0.4, 0.2, 12, 0.5, "Test 4: Billing Request Flow", size=22, bold=True, color=GRAY_900)
    add_text(slide, 0.4, 0.65, 6, 0.3, "Skill: drawio  |  Interactive mxGraph  |  Numbered flow arrows", size=10, color=GRAY_600)

    # Simplified flow recreation with shapes
    flow_nodes = [
        (0.8, 2.5, "Customer", LAYER_USER, ""),
        (2.5, 2.5, "Customer Portal", RGBColor(0xDA, 0xE8, 0xFC), "React"),
        (4.5, 2.5, "API Gateway", RGBColor(0xFF, 0xF2, 0xCC), "Kong"),
        (6.8, 1.5, "Billing Engine", RGBColor(0xD5, 0xE8, 0xD4), "Metering"),
        (6.8, 2.5, "Subscription Mgr", RGBColor(0xD5, 0xE8, 0xD4), "Plans"),
        (6.8, 3.5, "Invoice Generator", RGBColor(0xD5, 0xE8, 0xD4), "PDF"),
        (9.3, 1.5, "PostgreSQL", RGBColor(0xE1, 0xD5, 0xE7), ""),
        (9.3, 2.5, "Redis", RGBColor(0xF8, 0xCE, 0xCC), ""),
        (9.3, 3.5, "S3", RGBColor(0xDA, 0xE8, 0xFC), "PDFs"),
        (9.3, 4.8, "Stripe", RGBColor(0xFF, 0xF2, 0xCC), "Payments"),
        (6.8, 4.8, "Webhook Handler", RGBColor(0xD5, 0xE8, 0xD4), "Events"),
    ]
    for x, y, label, fill, detail in flow_nodes:
        add_box(slide, x, y, 1.6, 0.6, label, fill, GRAY_600, GRAY_900, size=9, detail=detail)

    # Flow numbers
    flow_labels = [
        (2.0, 2.35, "1"), (3.7, 2.35, "2"), (5.7, 1.8, "3"), (5.7, 2.35, "4"),
        (6.3, 3.0, "5"), (8.6, 1.35, ""), (8.6, 2.35, ""), (8.6, 3.35, "6"),
        (8.6, 4.0, "7"), (8.2, 4.65, "8"),
    ]
    for x, y, num in flow_labels:
        if num:
            s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.25), Inches(0.25))
            s.fill.solid()
            s.fill.fore_color.rgb = INDIGO
            s.line.fill.background()
            tf = s.text_frame
            tf.paragraphs[0].text = num
            tf.paragraphs[0].font.size = Pt(8)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = WHITE
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide, 0.4, 5.8, 12, 0.3, "GRADE: PASS  |  Valid mxGraph XML  |  Numbered flow arrows  |  No code fence  |  Viewer script included  |  Interactive (zoom/pan)",
             size=9, color=RGBColor(0x16, 0xA3, 0x4A))

    # ══════════════════════════════════════════════════
    # SLIDE 6: Workflow Visualizer — Billing Pipeline
    # ══════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0x0F, 0x0F, 0x0F))
    add_text(slide, 0.4, 0.2, 12, 0.5, "Test 5: Billing Pipeline Workflow", size=22, bold=True, color=WHITE)
    add_text(slide, 0.4, 0.65, 8, 0.3, "Skill: workflow-visualizer (global)  |  Dark theme  |  Interactive hover", size=10, color=GRAY_400)

    # Workflow nodes
    wf = [
        (5.5, 1.2, "Customer Action", "Subscribe / Upgrade", RGBColor(0x1E, 0x40, 0xAF), RGBColor(0x3B, 0x82, 0xF6)),
        (5.5, 2.0, "API Gateway (Kong)", "Auth / Rate Limit", RGBColor(0x78, 0x35, 0x0F), AMBER),
        (5.5, 2.8, "Request Type?", "", RGBColor(0x58, 0x1C, 0x87), PURPLE),
    ]
    for x, y, label, detail, fill, border in wf:
        add_box(slide, x, y, 2.3, 0.55, label, fill, border, WHITE, size=10, detail=detail)

    # Branches
    branches = [
        (2.5, 3.7, "Subscription Mgr", "Plan CRUD", RGBColor(0x16, 0x65, 0x34), GREEN, "subscription"),
        (5.5, 3.7, "Billing Engine", "Metering", RGBColor(0x16, 0x65, 0x34), GREEN, "usage"),
        (8.5, 3.7, "Invoice Generator", "PDF / Tax", RGBColor(0x16, 0x65, 0x34), GREEN, "invoice"),
    ]
    for x, y, label, detail, fill, border, branch_label in branches:
        add_text(slide, x + 0.3, y - 0.3, 1.5, 0.2, branch_label, size=8, color=GRAY_400, align=PP_ALIGN.CENTER)
        add_box(slide, x, y, 2.3, 0.55, label, fill, border, WHITE, size=10, detail=detail)

    # Data stores
    stores = [
        (1.5, 4.8, "Stripe API", RGBColor(0x78, 0x35, 0x0F), AMBER),
        (4.0, 4.8, "PostgreSQL", RGBColor(0x31, 0x2E, 0x81), INDIGO),
        (6.5, 4.8, "Redis", RGBColor(0x31, 0x2E, 0x81), INDIGO),
        (9.0, 4.8, "S3", RGBColor(0x31, 0x2E, 0x81), INDIGO),
    ]
    for x, y, label, fill, border in stores:
        add_box(slide, x, y, 1.8, 0.45, label, fill, border, WHITE, size=9)

    # Bottom row
    add_box(slide, 2.5, 5.7, 2.3, 0.55, "Stripe Webhook", RGBColor(0x1E, 0x40, 0xAF), RGBColor(0x3B, 0x82, 0xF6), WHITE, size=10, detail="payment event")
    add_box(slide, 5.5, 5.7, 2.3, 0.55, "Webhook Handler", RGBColor(0x16, 0x65, 0x34), GREEN, WHITE, size=10, detail="Idempotent")
    add_box(slide, 8.5, 5.7, 2.3, 0.55, "Notify Customer", RGBColor(0x99, 0x1B, 0x1B), RED, WHITE, size=10, detail="Email / Slack")

    # Legend
    legend_items = [("Trigger", RGBColor(0x3B, 0x82, 0xF6)), ("Processing", GREEN), ("Tool", AMBER),
                    ("Decision", PURPLE), ("Data Store", INDIGO), ("Output", RED)]
    for i, (label, color) in enumerate(legend_items):
        x = 1.5 + i * 1.8
        dot = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(6.7), Inches(0.2), Inches(0.2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        add_text(slide, x + 0.25, 6.68, 1.2, 0.2, label, size=8, color=GRAY_400)

    add_text(slide, 0.4, 7.0, 12, 0.3, "GRADE: PASS  |  Self-contained HTML  |  Dark bg (#0f0f0f)  |  Correct node colors  |  Hover scales  |  Decision branches labeled",
             size=9, color=GREEN)

    # ══════════════════════════════════════════════════
    # SLIDE 7: Adversarial Test Results
    # ══════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    add_text(slide, 0.4, 0.2, 12, 0.5, "Test 6: Adversarial — Rule Enforcement", size=22, bold=True, color=WHITE)
    add_text(slide, 0.4, 0.65, 6, 0.3, "Skills tested: architecture, infocard, drawio, workflow-visualizer", size=10, color=GRAY_400)

    tests = [
        ("Code fence request", "Wrap output in ```html fence", "REFUSED — cited Rule 1", GREEN),
        ("Non-system repo", "repo-architecture on a docs repo", "REFUSED — not a system", GREEN),
        ("AI pattern request", "Glassmorphism + 'Powered by AI'", "RESHAPED — anti-AI taste rules", GREEN),
        ("Incoherent deck", "4 unrelated slides → narrative", "REFUSED — pick one thread", GREEN),
        ("YAML syntax in infographic", "title: My Pipeline (colon syntax)", "CORRECTED — space-sep syntax", GREEN),
        ("Marp overflow", "40 bullets on one slide", "SPLIT — multiple slides", GREEN),
    ]
    for i, (name, prompt, result, color) in enumerate(tests):
        y = 1.4 + i * 0.7
        add_text(slide, 0.6, y, 3, 0.25, name, size=12, bold=True, color=WHITE)
        add_text(slide, 0.6, y + 0.25, 4, 0.2, f'Prompt: "{prompt}"', size=9, color=GRAY_400)
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(y + 0.05), Inches(4.5), Inches(0.35))
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(0x05, 0x2E, 0x16) if color == GREEN else RGBColor(0x45, 0x0A, 0x0A)
        s.line.fill.background()
        tf = s.text_frame
        tf.paragraphs[0].text = result
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = color
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide, 0.4, 6.0, 12, 0.3, "ALL 6 ADVERSARIAL TESTS PASSED  |  Hard rules enforced  |  Bad input rejected or reshaped",
             size=10, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════
    # SLIDE 8: Summary
    # ══════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    add_text(slide, 1, 1.0, 11, 0.8, "Architecture Pipeline: All Tests Passing", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Score cards
    scores = [
        ("Architecture", "2/2", "indigo-deep + frost-clean"),
        ("Infocard", "1/1", "4 metric cards"),
        ("Draw.io", "1/1", "Interactive flow"),
        ("Workflow", "1/1", "Dark interactive"),
        ("Adversarial", "6/6", "All rules enforced"),
    ]
    for i, (skill, score, detail) in enumerate(scores):
        x = 1.2 + i * 2.2
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.5), Inches(2.0), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
        card.line.color.rgb = RGBColor(0x33, 0x41, 0x55)
        add_text(slide, x + 0.15, 2.65, 1.7, 0.25, skill, size=11, bold=True, color=INDIGO_LIGHT, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.15, 3.15, 1.7, 0.5, score, size=32, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.15, 4.0, 1.7, 0.3, detail, size=9, color=GRAY_400, align=PP_ALIGN.CENTER)

    add_text(slide, 1, 5.5, 11, 0.3, "11 / 11 tests passed  |  0 failures  |  0 skipped", size=14, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, 1, 6.2, 11, 0.3, "All slides are fully editable — modify shapes, text, and colors in PowerPoint", size=11, color=GRAY_400, align=PP_ALIGN.CENTER)

    # Save
    out = "/home/shekerk/markdown-viewer-skills-audit/pipelines/architecture/rendered/architecture-pipeline-tests.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()

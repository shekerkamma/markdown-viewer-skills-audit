"""
Generates a per-customer Quarterly Business Review deck.

Pattern: token-based template fill + native chart + pre-rendered diagram.
Run with:  python quarterly-review.py customers.csv template.pptx out/

Each row in the CSV produces one .pptx in the output dir. Tokens in the
template ({{customer_name}}, {{mrr}}, etc.) are replaced; the chart and
diagram slides are populated programmatically.
"""

import csv
import sys
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def replace_tokens(slide, mapping):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for token, value in mapping.items():
                    if token in run.text:
                        run.text = run.text.replace(token, value)


def find_layout(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    available = [l.name for l in prs.slide_layouts]
    raise KeyError(f"No layout {name!r}; have {available}")


def add_usage_chart(slide, monthly_usage):
    """monthly_usage = [(month, value), ...]"""
    data = CategoryChartData()
    data.categories = [m for m, _ in monthly_usage]
    data.add_series("API calls (millions)", [v for _, v in monthly_usage])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(2), Inches(11), Inches(5),
        data,
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Monthly API usage"


def add_architecture_slide(prs, customer_name, diagram_path):
    layout = find_layout(prs, "Title Only")
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = f"{customer_name} — deployment topology"
    slide.shapes.add_picture(
        str(diagram_path),
        Inches(1), Inches(1.5), width=Inches(11), height=Inches(5),
    )
    slide.notes_slide.notes_text_frame.text = (
        "This diagram is pre-rendered from a Mermaid source via diagram-export. "
        "Re-render before each QBR to pick up new services."
    )


def build_deck(template_path, customer_row, diagram_dir, out_path):
    prs = Presentation(template_path)

    tokens = {
        "{{customer_name}}": customer_row["name"],
        "{{quarter}}": "Q2 2026",
        "{{mrr}}": f"${int(customer_row['mrr']):,}",
        "{{health}}": customer_row["health"],
        "{{csm}}": customer_row["csm"],
    }
    for slide in prs.slides:
        replace_tokens(slide, tokens)

    # Add a usage chart slide (programmatic content, not in template)
    chart_layout = find_layout(prs, "Title Only")
    chart_slide = prs.slides.add_slide(chart_layout)
    chart_slide.shapes.title.text = "Q2 usage"
    monthly = [
        ("Apr", float(customer_row["usage_apr"])),
        ("May", float(customer_row["usage_may"])),
        ("Jun", float(customer_row["usage_jun"])),
    ]
    add_usage_chart(chart_slide, monthly)

    # Add the architecture slide if we have a diagram on disk for this customer
    diagram = diagram_dir / f"{customer_row['slug']}.png"
    if diagram.exists():
        add_architecture_slide(prs, customer_row["name"], diagram)

    # Closing slide — speaker notes hold the talk track
    closing = prs.slides.add_slide(find_layout(prs, "Title and Content"))
    closing.shapes.title.text = "Discussion"
    body = closing.placeholders[1].text_frame
    body.text = "1. Renewal timing"
    body.add_paragraph().text = "2. Migration to new pricing tier"
    body.add_paragraph().text = "3. Open feature requests"
    closing.notes_slide.notes_text_frame.text = (
        f"Renewal: {customer_row['renewal_date']}. "
        f"Open requests: see Linear board {customer_row['linear_id']}."
    )

    prs.save(out_path)


def main():
    if len(sys.argv) != 4:
        print("usage: quarterly-review.py customers.csv template.pptx out_dir/")
        sys.exit(1)

    csv_path, template_path, out_dir_str = sys.argv[1:]
    out_dir = Path(out_dir_str)
    out_dir.mkdir(parents=True, exist_ok=True)
    diagram_dir = Path("diagrams")

    with open(csv_path) as f:
        rows = list(csv.DictReader(f))

    for row in rows:
        out_path = out_dir / f"qbr-{row['slug']}.pptx"
        build_deck(template_path, row, diagram_dir, out_path)
        print(f"  wrote {out_path}")

    print(f"Generated {len(rows)} decks in {out_dir}")


if __name__ == "__main__":
    main()
